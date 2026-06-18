"""Template analysis for MemSlides templates.

This module extracts slide HTML/text/image metadata and builds a deterministic
``slide_induction`` profile without relying on any legacy template engine.
"""

from __future__ import annotations

import hashlib
import html
import json
import logging
import re
import shutil
from collections import Counter, defaultdict
from dataclasses import dataclass, field
from datetime import datetime
from pathlib import Path
from typing import Any

from memslides.memory.core.template_models import (
    CanonicalLayout,
    CanonicalSlot,
    TemplateSemanticModel,
)
from memslides.templates.shell import (
    LAYERED_SHELL,
    NATIVE_BACKGROUND,
    NO_SHELL,
    RASTER_SHELL,
)

logger = logging.getLogger(__name__)


def _as_int(value: Any, default: int = 0) -> int:
    try:
        return int(value)
    except Exception:
        return default


def _emu_to_float(value: Any) -> float:
    try:
        return float(value)
    except Exception:
        return 0.0


def _color_to_hex(color: Any) -> str:
    rgb = getattr(color, "rgb", None)
    if rgb is None:
        return ""
    try:
        return f"#{int(rgb[0]):02x}{int(rgb[1]):02x}{int(rgb[2]):02x}"
    except Exception:
        return str(rgb)


def _slide_font_name(shape: Any) -> str:
    text_frame = getattr(shape, "text_frame", None)
    if text_frame is None:
        return ""
    for paragraph in getattr(text_frame, "paragraphs", []) or []:
        for run in getattr(paragraph, "runs", []) or []:
            font = getattr(run, "font", None)
            name = getattr(font, "name", "") if font is not None else ""
            if name:
                return str(name)
    return ""


def _slide_text_size(shape: Any) -> int:
    text_frame = getattr(shape, "text_frame", None)
    if text_frame is None:
        return 0
    sizes: list[int] = []
    for paragraph in getattr(text_frame, "paragraphs", []) or []:
        for run in getattr(paragraph, "runs", []) or []:
            font = getattr(run, "font", None)
            size = getattr(font, "size", None) if font is not None else None
            if size is not None:
                sizes.append(_as_int(round(_emu_to_float(size) / 12700.0)))
    if not sizes:
        return 0
    sizes.sort()
    return sizes[len(sizes) // 2]


def _shape_text(shape: Any) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    text = getattr(shape, "text", "") or ""
    return str(text).strip()


def _shape_is_image(shape: Any) -> bool:
    try:
        if getattr(shape, "image", None) is not None:
            return True
    except Exception:
        pass
    shape_type = str(getattr(shape, "shape_type", "")).lower()
    return "picture" in shape_type or "image" in shape_type


def _shape_type_label(shape: Any) -> str:
    return str(getattr(shape, "shape_type", "") or "").strip()


def _placeholder_type(shape: Any) -> str:
    try:
        return str(shape.placeholder_format.type)
    except Exception:
        return ""


def _placeholder_idx(shape: Any) -> int | None:
    try:
        return int(shape.placeholder_format.idx)
    except Exception:
        return None


def _is_placeholder(shape: Any) -> bool:
    if bool(getattr(shape, "is_placeholder", False)):
        return True
    return bool(_placeholder_type(shape))


def _is_ignored_placeholder(shape: Any) -> bool:
    placeholder_type = _placeholder_type(shape).upper()
    return any(
        token in placeholder_type
        for token in ("DATE", "FOOTER", "SLIDE_NUMBER", "HEADER")
    )


def _text_role_from_shape(
    *,
    shape: Any,
    text: str,
    shape_idx: int,
    left_pct: float,
    top_pct: float,
    width_pct: float,
    height_pct: float,
    font_size: int,
) -> str:
    placeholder_type = _placeholder_type(shape).upper()
    lowered = text.lower().strip()
    if _is_ignored_placeholder(shape):
        return "decorative"
    if any(token in lowered for token in ("汇报人", "答辩人", "报告人", "时间：", "日期：", "beijing university", "bupt")):
        return "decorative"
    if re.fullmatch(r"[\d\s./:-]+", text.strip()) and len(text.strip()) <= 12:
        return "number_badge"
    if len(text.strip()) <= 2 and re.fullmatch(r"[\u4e00-\u9fffA-Za-z]{1,2}", text.strip()):
        return "decorative"
    if "/" in text and len(text.strip()) <= 24:
        return "section_label"
    if any(token in lowered for token in ("please enter a title", "add text", "输入标题", "添加小标题", "关键词")):
        if top_pct <= 24:
            return "title"
        return "section_label"
    if any(token in lowered for token in ("contents", "目录", "agenda", "outline")):
        return "title"
    if any(token in lowered for token in ("background and significance", "research methods", "applications")):
        return "toc_item"
    if "SUBTITLE" in placeholder_type:
        return "subtitle"
    if "TITLE" in placeholder_type or "CENTER_TITLE" in placeholder_type:
        return "title"
    if "BODY" in placeholder_type or "OBJECT" in placeholder_type:
        return "body"
    if "CAPTION" in placeholder_type:
        return "caption"
    if top_pct <= 22 and width_pct <= 18 and len(text) <= 12:
        return "number_badge"
    if top_pct <= 22 and height_pct <= 8 and len(text) <= 32:
        return "section_label"
    if top_pct <= 20 and (font_size >= 24 or len(text) <= 80):
        return "title"
    if font_size >= 28 and len(text) <= 120:
        return "title"
    if height_pct <= 8 and len(text) <= 36:
        return "label"
    if "\n" in text or len(text) > 80:
        return "body"
    if shape_idx == 0 and len(text) <= 120:
        return "title"
    return "body"


def _text_kind_from_role(role: str) -> str:
    if role == "title":
        return "h1"
    if role == "subtitle":
        return "h2"
    if role in {"body", "caption"}:
        return "p"
    return "span"


def _iter_shape_tree(
    shapes: Any,
    *,
    parent_left: float = 0.0,
    parent_top: float = 0.0,
    path_prefix: str = "",
    depth: int = 0,
):
    """Yield shapes recursively with approximate absolute coordinates.

    PowerPoint templates often store real text boxes inside GROUP shapes.  The
    previous analyzer only inspected top-level slide shapes, so those templates
    were misread as having no text slots.  python-pptx exposes group children in
    group-local coordinates; adding parent offsets is not a full affine
    transform, but it is stable enough for slot discovery and percent geometry.
    """

    for idx, shape in enumerate(list(shapes or [])):
        left = parent_left + _emu_to_float(getattr(shape, "left", 0))
        top = parent_top + _emu_to_float(getattr(shape, "top", 0))
        path = f"{path_prefix}.{idx}" if path_prefix else str(idx)
        yield shape, path, left, top, depth, idx
        child_shapes = getattr(shape, "shapes", None)
        if child_shapes is not None:
            yield from _iter_shape_tree(
                child_shapes,
                parent_left=left,
                parent_top=top,
                path_prefix=path,
                depth=depth + 1,
            )


def _classify_image_semantics(
    left_pct: float,
    top_pct: float,
    width_pct: float,
    height_pct: float,
) -> str:
    area = max(0.0, width_pct) * max(0.0, height_pct)
    if area >= 70 or (width_pct >= 88 and height_pct >= 60):
        return "background"
    if area <= 4 or height_pct <= 8 or width_pct <= 8:
        return "decorative"
    if top_pct <= 8 and (left_pct <= 12 or left_pct >= 75):
        return "decorative"
    if top_pct >= 84 or left_pct >= 88:
        return "decorative"
    return "content"


def _region_from_shape(shape: dict[str, Any], *, role: str = "") -> dict[str, Any]:
    return {
        "role": role or str(shape.get("semantic_role", "") or "region"),
        "left_pct": float(shape.get("left_pct", 0) or 0),
        "top_pct": float(shape.get("top_pct", 0) or 0),
        "width_pct": float(shape.get("width_pct", 0) or 0),
        "height_pct": float(shape.get("height_pct", 0) or 0),
    }


def _shape_area_pct(shape: dict[str, Any]) -> float:
    return max(0.0, float(shape.get("width_pct", 0) or 0)) * max(
        0.0, float(shape.get("height_pct", 0) or 0)
    )


def _shape_shell_layer(shape: dict[str, Any], *, layer_type: str, z_index: int) -> dict[str, Any]:
    payload = {
        "type": layer_type,
        "role": str(shape.get("semantic_role", "") or layer_type),
        "z_index": z_index,
        "geometry": _region_from_shape(shape, role=str(shape.get("semantic_role", "") or layer_type)),
        "position": str(shape.get("position", "") or ""),
        "shape_path": str(shape.get("shape_path", "") or ""),
        "shape_type": str(shape.get("shape_type", "") or ""),
    }
    if shape.get("asset_name"):
        payload["asset_name"] = str(shape.get("asset_name", "") or "")
    if shape.get("text"):
        payload["text"] = str(shape.get("text", "") or "")
    return payload


def _protected_region_from_decorative_shape(shape: dict[str, Any]) -> dict[str, Any] | None:
    role = str(shape.get("semantic_role", "") or "")
    text = str(shape.get("text", "") or "").strip()
    area = _shape_area_pct(shape)
    top_pct = float(shape.get("top_pct", 0) or 0)
    height_pct = float(shape.get("height_pct", 0) or 0)
    if role == "background":
        return None
    if top_pct <= 18 and height_pct <= 18:
        region_role = "title_band"
    elif role in {"number_badge", "decorative"} and area <= 180:
        region_role = role
    elif text:
        region_role = "decorative_text"
    else:
        region_role = "decorative"
    return _region_from_shape(shape, role=region_role)


def _shape_position_label(left_pct: float, top_pct: float, width_pct: float, height_pct: float) -> str:
    if top_pct < 18:
        vertical = "top"
    elif top_pct > 66:
        vertical = "bottom"
    else:
        vertical = "middle"
    if left_pct < 20:
        horizontal = "left"
    elif left_pct > 65:
        horizontal = "right"
    else:
        horizontal = "center"
    if vertical == "middle" and horizontal == "center":
        return "center"
    return f"{vertical}-{horizontal}"


def _surface_mode_from_signature(signature: dict[str, Any]) -> str:
    if signature.get("supports_figure") or signature.get("supports_table") or signature.get("supports_chart"):
        return "figure_panel"
    if signature.get("supports_long_body") or signature.get("body_slot_count", 0) >= 1:
        return "light_surface"
    return "shell_text"


def _allowed_text_roles_from_slots(slots: list[dict[str, Any]]) -> list[str]:
    roles = []
    for slot in slots:
        role = str(slot.get("semantic_role", "") or "").strip()
        if role and role not in roles and role != "decorative":
            roles.append(role)
    return roles


def _allowed_visual_roles_from_slots(slots: list[dict[str, Any]]) -> list[str]:
    roles = []
    for slot in slots:
        role = str(slot.get("semantic_role", "") or "").strip()
        slot_type = str(slot.get("type", "") or "").strip()
        if slot_type in {"image", "chart", "table"}:
            normalized = role or slot_type
            if normalized not in roles:
                roles.append(normalized)
    return roles


def _detect_language(texts: list[str]) -> dict[str, str]:
    payload = "\n".join(texts[:6]).strip()
    if not payload:
        return {"lid": "unknown"}
    cjk = sum(1 for ch in payload if "\u4e00" <= ch <= "\u9fff")
    latin = sum(1 for ch in payload if ch.isascii() and ch.isalpha())
    if cjk > latin:
        return {"lid": "zh"}
    if latin > 0:
        return {"lid": "en"}
    return {"lid": "unknown"}


def _clean_text(text: str) -> str:
    text = re.sub(r"\s+", " ", text or "").strip()
    return html.escape(text)


def _sanitize_layout_token(text: str) -> str:
    cleaned = re.sub(r"[^a-zA-Z0-9]+", "_", str(text or "").strip().lower()).strip("_")
    return cleaned or "layout"


def _capacity_hint_from_slot(slot: dict[str, Any]) -> dict[str, Any]:
    sample_values = [str(item or "") for item in slot.get("data", []) or [] if str(item or "").strip()]
    font_size = _as_int(slot.get("font_size", 0), 0)
    width_pct = float(slot.get("width_pct", 0) or 0)
    height_pct = float(slot.get("height_pct", 0) or 0)
    baseline = max((len(value) for value in sample_values), default=0)
    if baseline <= 0 and slot.get("type") == "text":
        area = max(width_pct * height_pct, 1.0)
        if font_size >= 28:
            baseline = 36 if area < 250 else 72
        elif font_size >= 20:
            baseline = 80 if area < 500 else 140
        else:
            baseline = 120 if area < 500 else 220
    line_estimate = 1
    if font_size > 0 and height_pct > 0:
        line_estimate = max(1, int(round((height_pct * 7.2) / max(font_size, 12))))
    return {
        "suggested_max_chars": baseline,
        "sample_char_count": max((len(value) for value in sample_values), default=0),
        "line_capacity_estimate": line_estimate,
        "font_size": font_size,
    }


def _layout_archetype_from_slots(slots: list[dict[str, Any]]) -> str:
    roles = [str(slot.get("semantic_role", "") or "") for slot in slots if slot.get("type") == "text"]
    image_count = sum(1 for slot in slots if slot.get("type") == "image")
    body_count = sum(1 for role in roles if role == "body")
    title_count = sum(1 for role in roles if role == "title")
    caption_count = sum(1 for role in roles if role == "caption")
    if title_count and not body_count and image_count == 0:
        return "title_page"
    if title_count and body_count and image_count:
        return "title_content_visual"
    if title_count and body_count >= 2 and image_count == 0:
        return "title_content"
    if body_count >= 2 and image_count >= 2:
        return "multi_panel"
    if image_count >= 2 and body_count <= 1:
        return "picture_grid"
    if image_count == 1 and caption_count:
        return "picture_caption"
    if image_count == 1 and body_count:
        return "figure_explanation"
    if body_count >= 3:
        return "outline"
    return "content"


def _visual_pattern_from_slots(slots: list[dict[str, Any]]) -> str:
    text_count = sum(1 for slot in slots if slot.get("type") == "text")
    image_count = sum(1 for slot in slots if slot.get("type") == "image")
    roles = [str(slot.get("semantic_role", "") or "") for slot in slots if slot.get("type") == "text"]
    title_count = sum(1 for role in roles if role == "title")
    section_count = sum(1 for role in roles if role == "section_label")
    body_count = sum(1 for role in roles if role == "body")
    leftish = sum(
        1
        for slot in slots
        if float(slot.get("left_pct", 0) or 0) <= 18 and float(slot.get("height_pct", 0) or 0) >= 14
    )
    topish = sum(
        1
        for slot in slots
        if float(slot.get("top_pct", 0) or 0) <= 22 and float(slot.get("width_pct", 0) or 0) >= 25
    )
    if section_count >= 2 and body_count == 0 and image_count == 0:
        return "section_divider"
    if leftish >= 1 and title_count >= 1 and body_count >= 1:
        return "vertical_title"
    if topish >= 1 and body_count >= 1 and image_count == 0:
        return "top_title_text"
    if image_count == 0 and text_count <= 2:
        return "hero_text"
    if image_count == 0 and text_count >= 4:
        return "dense_text"
    if image_count == 1 and text_count >= 2:
        return "text_plus_visual"
    if image_count >= 2 and text_count <= 1:
        return "gallery"
    if image_count >= 1 and text_count >= 3:
        return "mixed_grid"
    return "balanced"


def _is_likely_decorative_text_shape(shape: dict[str, Any]) -> bool:
    role = str(shape.get("semantic_role", "") or "")
    text = str(shape.get("text", "") or "").strip()
    width_pct = float(shape.get("width_pct", 0) or 0)
    height_pct = float(shape.get("height_pct", 0) or 0)
    top_pct = float(shape.get("top_pct", 0) or 0)
    if role == "decorative":
        return True
    if role == "number_badge" and len(text) <= 4:
        return True
    if len(text) <= 2 and height_pct >= 16:
        return True
    if top_pct >= 84 and len(text) <= 24:
        return True
    if width_pct <= 8 and height_pct >= 14 and len(text) <= 3:
        return True
    if any(token in text.lower() for token in ("bupt", "beijing university")):
        return True
    return False


def _layout_signature_from_slots(slots: list[CanonicalSlot]) -> tuple:
    signature = []
    for slot in sorted(
        slots,
        key=lambda item: (
            item.role,
            item.type,
            round(float((item.geometry or {}).get("left_pct", 0) or 0), 1),
            round(float((item.geometry or {}).get("top_pct", 0) or 0), 1),
        ),
    ):
        geo = slot.geometry or {}
        signature.append(
            (
                slot.role,
                slot.type,
                round(float(geo.get("left_pct", 0) or 0) / 5) * 5,
                round(float(geo.get("top_pct", 0) or 0) / 5) * 5,
                round(float(geo.get("width_pct", 0) or 0) / 5) * 5,
                round(float(geo.get("height_pct", 0) or 0) / 5) * 5,
            )
        )
    return tuple(signature)


def _capability_signature_from_slots(slots: list[dict[str, Any]]) -> dict[str, Any]:
    text_slots = [slot for slot in slots if slot.get("type") == "text"]
    non_text_slots = [slot for slot in slots if slot.get("type") in {"image", "chart", "table"}]
    roles = [str(slot.get("semantic_role", "") or "") for slot in text_slots]
    body_slots = [slot for slot in text_slots if str(slot.get("semantic_role", "") or "") == "body"]
    widths = [float(slot.get("width_pct", 0) or 0) for slot in body_slots]
    top_positions = [float(slot.get("top_pct", 0) or 0) for slot in text_slots]
    left_positions = [float(slot.get("left_pct", 0) or 0) for slot in text_slots]
    image_roles = [str(slot.get("semantic_role", "") or "") for slot in non_text_slots]
    image_count = sum(1 for slot in non_text_slots if slot.get("type") == "image")
    chart_count = sum(1 for slot in non_text_slots if slot.get("type") == "chart")
    table_count = sum(1 for slot in non_text_slots if slot.get("type") == "table")
    title_count = sum(1 for role in roles if role == "title")
    section_count = sum(1 for role in roles if role == "section_label")
    caption_count = sum(1 for role in roles if role == "caption")
    body_count = len(body_slots)
    vertical_title = any(
        float(slot.get("left_pct", 0) or 0) <= 18 and float(slot.get("height_pct", 0) or 0) >= 14
        for slot in text_slots
    )
    return {
        "supports_title": title_count > 0,
        "supports_subtitle": any(role == "subtitle" for role in roles),
        "supports_long_body": any(float(slot.get("height_pct", 0) or 0) >= 20 for slot in body_slots),
        "supports_short_body": body_count > 0,
        "supports_figure": image_count > 0 or "content" in image_roles,
        "supports_chart": chart_count > 0 or "chart" in image_roles,
        "supports_table": table_count > 0 or "table" in image_roles,
        "supports_multi_column": sum(1 for width in widths if width <= 45) >= 2,
        "supports_section_divider": section_count >= 2 and body_count == 0,
        "supports_toc": section_count >= 4 and body_count == 0,
        "supports_ending": title_count >= 1 and body_count == 0 and image_count == 0 and caption_count == 0,
        "supports_caption": caption_count > 0,
        "vertical_title": vertical_title,
        "top_title": any(pos <= 20 for pos in top_positions) and title_count > 0,
        "left_title": any(pos <= 18 for pos in left_positions) and title_count > 0,
        "body_slot_count": body_count,
        "image_slot_count": image_count,
        "text_slot_count": len(text_slots),
    }


def _template_audit_payload(semantic_model: dict[str, Any]) -> dict[str, Any]:
    canonical_layouts = list(semantic_model.get("canonical_layouts", []) or [])
    template_shell = semantic_model.get("template_shell", {}) or {}
    return {
        "canonical_layout_count": len(canonical_layouts),
        "canonical_layouts": [
            {
                "name": layout.get("name", ""),
                "aliases": layout.get("aliases", []),
                "layout_archetype": layout.get("layout_archetype", ""),
                "visual_pattern": layout.get("visual_pattern", ""),
                "slot_roles": layout.get("slot_roles", []),
                "slot_count": layout.get("slot_count", 0),
                "slots": [
                    {
                        "name": slot.get("name", ""),
                        "type": slot.get("type", ""),
                        "role": slot.get("role", ""),
                        "geometry": slot.get("geometry", {}),
                        "capacity_hint": slot.get("capacity_hint", {}),
                    }
                    for slot in layout.get("slots", []) or []
                ],
            }
            for layout in canonical_layouts
        ],
        "template_shell": {
            "summary": template_shell.get("summary", {}),
            "layout_shells": [
                {
                    "layout": name,
                    "shell_id": data.get("shell_id", ""),
                    "shell_mode": data.get("shell_mode", ""),
                    "protected_regions": len(data.get("protected_regions", []) or []),
                    "background_assets": len(data.get("background_assets", []) or []),
                    "decorative_layers": len(data.get("decorative_layers", []) or []),
                }
                for name, data in (template_shell.get("layouts", {}) or {}).items()
                if isinstance(data, dict)
            ],
        },
        "style_system": semantic_model.get("style_system", {}),
        "decorative_system": semantic_model.get("decorative_system", {}),
        "suspicious_ambiguities": semantic_model.get("suspicious_ambiguities", []),
        "low_confidence": bool(semantic_model.get("low_confidence", False)),
    }


@dataclass
class TemplateAnalysis:
    template_path: str = ""
    template_name: str = ""
    slide_count: int = 0
    aspect_ratio: str = "16:9"
    slides_html: list[str] = field(default_factory=list)
    slides_text: list[str] = field(default_factory=list)
    slides_images: list[list[str]] = field(default_factory=list)
    layout_induction: dict[str, Any] = field(default_factory=dict)
    semantic_model: dict[str, Any] = field(default_factory=dict)
    image_stats: dict[str, Any] = field(default_factory=dict)
    shape_geometry: dict[str, Any] = field(default_factory=dict)
    font_names: list[str] = field(default_factory=list)
    total_slides: int = 0
    parsed_slides: int = 0
    skipped_slides: list[dict[str, Any]] = field(default_factory=list)
    created_at: str = field(default_factory=lambda: datetime.now().isoformat())

    def to_dict(self) -> dict:
        return {
            "template_path": self.template_path,
            "template_name": self.template_name,
            "slide_count": self.slide_count,
            "aspect_ratio": self.aspect_ratio,
            "slides_html": self.slides_html,
            "slides_text": self.slides_text,
            "slides_images": self.slides_images,
            "layout_induction": self.layout_induction,
            "semantic_model": self.semantic_model,
            "image_stats": self.image_stats,
            "shape_geometry": self.shape_geometry,
            "font_names": self.font_names,
            "total_slides": self.total_slides,
            "parsed_slides": self.parsed_slides,
            "skipped_slides": self.skipped_slides,
            "created_at": self.created_at,
        }


class TemplateAnalyzer:
    """Deterministic PPT template analyzer."""

    def __init__(
        self,
        workspace: Path | str | None = None,
        language_model: Any = None,
        vision_model: Any = None,
        **_: Any,
    ):
        self.workspace = Path(workspace) if workspace else Path.cwd()
        self.language_model = language_model
        self.vision_model = vision_model
        self._legacy_backend_path = None
        self._temp_dir = self.workspace / ".template_analysis"

    async def analyze(self, pptx_path: Path | str, output_dir: Path | str | None = None) -> TemplateAnalysis:
        pptx_path = Path(pptx_path)
        if not pptx_path.exists():
            raise FileNotFoundError(f"PPT file not found: {pptx_path}")

        try:
            from pptx import Presentation as PptxPresentation
        except ImportError as exc:  # pragma: no cover - dependency issue
            raise RuntimeError("python-pptx is required for template analysis") from exc

        temp_dir = Path(output_dir) if output_dir else self._temp_dir / pptx_path.stem
        temp_dir.mkdir(parents=True, exist_ok=True)
        images_dir = temp_dir / "images"
        images_dir.mkdir(parents=True, exist_ok=True)

        prs = PptxPresentation(str(pptx_path))
        slide_width = _emu_to_float(getattr(prs, "slide_width", 0)) or 1.0
        slide_height = _emu_to_float(getattr(prs, "slide_height", 0)) or 1.0
        aspect_ratio = self._get_aspect_ratio(slide_width, slide_height)
        layout_capabilities = self._extract_layout_capabilities(
            prs=prs,
            slide_width=slide_width,
            slide_height=slide_height,
        )

        slide_records: list[dict[str, Any]] = []
        slides_html: list[str] = []
        slides_text: list[str] = []
        slides_images: list[list[str]] = []
        image_stats: dict[str, Any] = {}
        shape_geometry: dict[str, Any] = {}
        font_names: list[str] = []
        skipped_slides: list[dict[str, Any]] = []

        for slide_idx, slide in enumerate(prs.slides, start=1):
            try:
                record = self._extract_slide_record(
                    slide=slide,
                    slide_idx=slide_idx,
                    slide_width=slide_width,
                    slide_height=slide_height,
                    images_dir=images_dir,
                    image_stats=image_stats,
                    font_names=font_names,
                    layout_capabilities=layout_capabilities,
                )
                slide_records.append(record)
                slides_html.append(record["html"])
                slides_text.append(record["text"])
                slides_images.append(record["images"])
                shape_geometry[str(slide_idx - 1)] = record["shape_geometry"]
            except Exception as exc:
                logger.warning("Template slide %s skipped: %s", slide_idx, exc)
                skipped_slides.append({"slide_idx": slide_idx, "error": str(exc)})
                slides_html.append(f'<div class="slide" data-index="{slide_idx}"></div>')
                slides_text.append("")
                slides_images.append([])
                shape_geometry[str(slide_idx - 1)] = {
                    "by_text": {},
                    "by_order": [],
                    "text_shapes": [],
                    "image_shapes": [],
                }

        layout_induction = self._build_layout_induction(slide_records, prs, layout_capabilities)
        semantic_model = self._build_semantic_model(
            slide_records=slide_records,
            layout_induction=layout_induction,
            layout_capabilities=layout_capabilities,
            aspect_ratio=aspect_ratio,
        )
        total_slides = len(prs.slides)
        parsed_slides = len(slide_records)

        analysis = TemplateAnalysis(
            template_path=str(pptx_path),
            template_name=pptx_path.stem,
            slide_count=parsed_slides,
            aspect_ratio=aspect_ratio,
            slides_html=slides_html,
            slides_text=slides_text,
            slides_images=slides_images,
            layout_induction=layout_induction,
            semantic_model=semantic_model.to_dict(),
            image_stats=image_stats,
            shape_geometry=shape_geometry,
            font_names=self._most_common_fonts(font_names),
            total_slides=total_slides,
            parsed_slides=parsed_slides,
            skipped_slides=skipped_slides,
        )

        self._persist_analysis_outputs(temp_dir, analysis)
        return analysis

    def _extract_slide_record(
        self,
        *,
        slide: Any,
        slide_idx: int,
        slide_width: float,
        slide_height: float,
        images_dir: Path,
        image_stats: dict[str, Any],
        font_names: list[str],
        layout_capabilities: dict[str, dict[str, Any]],
    ) -> dict[str, Any]:
        text_fragments: list[str] = []
        html_parts: list[str] = [f'<div class="slide" data-index="{slide_idx}">']
        image_paths: list[str] = []
        by_order: list[dict[str, Any]] = []
        by_text: dict[str, dict[str, Any]] = {}
        text_shapes: list[dict[str, Any]] = []
        image_shapes: list[dict[str, Any]] = []

        shapes = list(_iter_shape_tree(getattr(slide, "shapes", []) or []))
        for shape_idx, (shape, shape_path, left, top, depth, sibling_idx) in enumerate(shapes):
            width = _emu_to_float(getattr(shape, "width", 0))
            height = _emu_to_float(getattr(shape, "height", 0))
            left_pct = round(left / slide_width * 100, 2)
            top_pct = round(top / slide_height * 100, 2)
            width_pct = round(width / slide_width * 100, 2)
            height_pct = round(height / slide_height * 100, 2)
            position = _shape_position_label(left_pct, top_pct, width_pct, height_pct)

            text = _shape_text(shape)
            font_name = _slide_font_name(shape)
            font_size = _slide_text_size(shape)
            if font_name:
                font_names.append(font_name)

            if _shape_is_image(shape):
                image_semantics = _classify_image_semantics(
                    left_pct, top_pct, width_pct, height_pct
                )
                image_info = getattr(shape, "image", None)
                stored_name = ""
                if image_info is not None:
                    stored_name = self._write_image_asset(image_info, images_dir, slide_idx, shape_idx)
                if not stored_name:
                    stored_name = f"image_{slide_idx}_{shape_idx + 1}"
                image_paths.append(stored_name)
                self._update_image_stats(
                    image_stats=image_stats,
                    image_name=stored_name,
                    slide_idx=slide_idx,
                    slide_width=slide_width,
                    slide_height=slide_height,
                    shape_width=width,
                    shape_height=height,
                )
                shape_record = {
                    "left_pct": left_pct,
                    "top_pct": top_pct,
                    "width_pct": width_pct,
                    "height_pct": height_pct,
                    "position": position,
                    "font_name": font_name,
                    "font_size": font_size,
                    "text": text,
                    "type": "image",
                    "asset_name": stored_name,
                    "semantic_role": image_semantics,
                    "is_content": image_semantics == "content",
                    "shape_path": shape_path,
                    "shape_type": _shape_type_label(shape),
                    "depth": depth,
                    "source": "sample",
                }
                by_order.append(shape_record)
                image_shapes.append(shape_record)
                html_parts.append(
                    f'<img data-name="{html.escape(stored_name)}" '
                    f'style="position:absolute;left:{left_pct:.2f}%;top:{top_pct:.2f}%;'
                    f'width:{width_pct:.2f}%;height:{height_pct:.2f}%;" '
                    f'src="{html.escape(stored_name)}" alt="{html.escape(stored_name)}" />'
                )
                continue

            if text:
                text_fragments.append(text)
                semantic_role = _text_role_from_shape(
                    shape=shape,
                    text=text,
                    shape_idx=shape_idx,
                    left_pct=left_pct,
                    top_pct=top_pct,
                    width_pct=width_pct,
                    height_pct=height_pct,
                    font_size=font_size,
                )
                if semantic_role == "decorative":
                    is_content_text = False
                else:
                    is_content_text = True
                text_kind = _text_kind_from_role(semantic_role)
                shape_record = {
                    "left_pct": left_pct,
                    "top_pct": top_pct,
                    "width_pct": width_pct,
                    "height_pct": height_pct,
                    "position": position,
                    "font_name": font_name,
                    "font_size": font_size,
                    "text": text,
                    "type": "text",
                    "semantic_role": semantic_role,
                    "is_content": is_content_text,
                    "shape_path": shape_path,
                    "shape_type": _shape_type_label(shape),
                    "depth": depth,
                    "source": "sample",
                    "placeholder_type": _placeholder_type(shape),
                    "placeholder_idx": _placeholder_idx(shape),
                }
                by_order.append(shape_record)
                text_shapes.append(shape_record)
                by_text[text] = shape_record
                style_bits = [
                    "position:absolute",
                    f"left:{left_pct:.2f}%",
                    f"top:{top_pct:.2f}%",
                    f"width:{width_pct:.2f}%",
                    f"height:{height_pct:.2f}%",
                ]
                if font_name:
                    style_bits.append(f"font-family:{font_name}")
                if font_size:
                    style_bits.append(f"font-size:{font_size}px")
                rendered_text = _clean_text(text).replace("\n", "<br/>")
                html_parts.append(
                    f'<div class="{text_kind}" style="{";".join(style_bits)}">{rendered_text}</div>'
                )

        html_parts.append("</div>")
        slide_layout_name = str(getattr(getattr(slide, "slide_layout", None), "name", "") or "default").strip()
        editable_slots = self._editable_slots_for_slide(slide, layout_capabilities)
        sample_elements = self._slide_elements(text_shapes, image_shapes, include_decorative=False)
        fallback_elements = sample_elements
        if not fallback_elements:
            fallback_elements = self._slide_elements(text_shapes, image_shapes, include_decorative=True)
        elements = self._merge_layout_slots_with_sample_elements(
            editable_slots=editable_slots,
            sample_elements=fallback_elements,
        )
        return {
            "index": slide_idx,
            "html": "\n".join(html_parts),
            "text": "\n".join(text_fragments),
            "images": image_paths,
            "shape_geometry": {
                "by_text": by_text,
                "by_order": by_order,
                "text_shapes": text_shapes,
                "image_shapes": image_shapes,
            },
            "layout_name": slide_layout_name or "default",
            "content_type": self._infer_content_type(text_fragments, image_paths),
            "slide_title": self._slide_title(text_fragments),
            "image_count": sum(
                1 for shape in image_shapes if shape.get("semantic_role") == "content"
            ),
            "text_count": sum(1 for shape in text_shapes if shape.get("is_content")),
            "sample_text_count": len(text_fragments),
            "editable_text_count": sum(1 for slot in editable_slots if slot.get("type") == "text"),
            "editable_image_count": sum(1 for slot in editable_slots if slot.get("type") == "image"),
            "elements": elements,
            "sample_elements": sample_elements,
            "editable_slots": editable_slots,
        }

    def _merge_layout_slots_with_sample_elements(
        self,
        *,
        editable_slots: list[dict[str, Any]],
        sample_elements: list[dict[str, Any]],
    ) -> list[dict[str, Any]]:
        if not editable_slots:
            return sample_elements
        if not sample_elements:
            return editable_slots

        slot_text_count = sum(1 for slot in editable_slots if slot.get("type") == "text")
        sample_text = [element for element in sample_elements if element.get("type") == "text"]
        sample_images = [element for element in sample_elements if element.get("type") == "image"]

        if slot_text_count == 0 and sample_text:
            merged = list(sample_text)
            slot_image_names = {str(slot.get("name", "")) for slot in editable_slots if slot.get("type") == "image"}
            merged.extend(editable_slots)
            merged.extend(
                element
                for element in sample_images
                if str(element.get("name", "")) not in slot_image_names
            )
            return merged

        return editable_slots

    def _slide_elements(
        self,
        text_shapes: list[dict[str, Any]],
        image_shapes: list[dict[str, Any]],
        *,
        include_decorative: bool = False,
    ) -> list[dict[str, Any]]:
        elements: list[dict[str, Any]] = []
        content_text_shapes = [
            shape for shape in text_shapes
            if (include_decorative or shape.get("is_content")) and not _is_likely_decorative_text_shape(shape)
        ]
        role_counts: Counter[str] = Counter()
        for idx, shape in enumerate(content_text_shapes):
            text = str(shape.get("text", "") or "")
            role = str(shape.get("semantic_role", "") or "text")
            role_counts[role] += 1
            if role_counts[role] == 1:
                label = role if role in {"title", "subtitle", "body", "caption"} else f"text_{idx + 1}"
            else:
                label = f"{role}_{role_counts[role]}"
            elements.append(
                {
                    "name": label,
                    "type": "text",
                    "data": [text],
                    "semantic_role": role,
                    "source": shape.get("source", "sample"),
                    "left_pct": shape.get("left_pct", 0),
                    "top_pct": shape.get("top_pct", 0),
                    "width_pct": shape.get("width_pct", 0),
                    "height_pct": shape.get("height_pct", 0),
                }
            )
        content_images = [
            shape for shape in image_shapes
            if include_decorative or shape.get("semantic_role") == "content"
        ]
        for idx, shape in enumerate(content_images):
            image_name = shape.get("asset_name") or f"content_image_{idx + 1}"
            elements.append(
                {
                    "name": f"image_{idx + 1}",
                    "type": "image",
                    "data": [image_name],
                    "is_content": shape.get("semantic_role") == "content",
                    "semantic_role": shape.get("semantic_role", "content"),
                    "source": shape.get("source", "sample"),
                    "left_pct": shape.get("left_pct", 0),
                    "top_pct": shape.get("top_pct", 0),
                    "width_pct": shape.get("width_pct", 0),
                    "height_pct": shape.get("height_pct", 0),
                }
            )
        return elements

    def _slide_title(self, text_fragments: list[str]) -> str:
        return text_fragments[0][:120] if text_fragments else ""

    def _extract_layout_capabilities(
        self,
        *,
        prs: Any,
        slide_width: float,
        slide_height: float,
    ) -> dict[str, dict[str, Any]]:
        capabilities: dict[str, dict[str, Any]] = {}
        layouts = list(getattr(prs, "slide_layouts", []) or [])
        for layout_idx, layout in enumerate(layouts):
            layout_name = str(getattr(layout, "name", "") or f"layout_{layout_idx}").strip()
            slots: list[dict[str, Any]] = []
            for shape_idx, shape in enumerate(list(getattr(layout, "shapes", []) or [])):
                if not _is_placeholder(shape) or _is_ignored_placeholder(shape):
                    continue
                slot = self._slot_from_placeholder(
                    shape=shape,
                    shape_idx=shape_idx,
                    slide_width=slide_width,
                    slide_height=slide_height,
                    source="layout",
                )
                if slot:
                    slots.append(slot)
            capabilities[layout_name or f"layout_{layout_idx}"] = {
                "layout_index": layout_idx,
                "layout_name": layout_name,
                "editable_slots": slots,
                "text_slot_count": sum(1 for slot in slots if slot.get("type") == "text"),
                "image_slot_count": sum(1 for slot in slots if slot.get("type") == "image"),
            }
        return capabilities

    def _slot_from_placeholder(
        self,
        *,
        shape: Any,
        shape_idx: int,
        slide_width: float,
        slide_height: float,
        source: str,
    ) -> dict[str, Any] | None:
        left = _emu_to_float(getattr(shape, "left", 0))
        top = _emu_to_float(getattr(shape, "top", 0))
        width = _emu_to_float(getattr(shape, "width", 0))
        height = _emu_to_float(getattr(shape, "height", 0))
        left_pct = round(left / slide_width * 100, 2)
        top_pct = round(top / slide_height * 100, 2)
        width_pct = round(width / slide_width * 100, 2)
        height_pct = round(height / slide_height * 100, 2)
        placeholder_type = _placeholder_type(shape)
        placeholder_type_upper = placeholder_type.upper()
        raw_name = str(getattr(shape, "name", "") or "").strip()
        sample_text = _shape_text(shape)
        font_name = _slide_font_name(shape)
        font_size = _slide_text_size(shape)

        if "PICTURE" in placeholder_type_upper:
            slot_type = "image"
            role = "image"
        elif "OBJECT" in placeholder_type_upper:
            slot_type = "text"
            role = "body"
        else:
            slot_type = "text"
            role = _text_role_from_shape(
                shape=shape,
                text=sample_text or raw_name,
                shape_idx=shape_idx,
                left_pct=left_pct,
                top_pct=top_pct,
                width_pct=width_pct,
                height_pct=height_pct,
                font_size=font_size,
            )
            if role == "decorative":
                return None

        name_base = role if role in {"title", "subtitle", "body", "caption", "image"} else "slot"
        placeholder_idx = _placeholder_idx(shape)
        suffix = placeholder_idx if placeholder_idx is not None else shape_idx + 1
        name = name_base if suffix in (0, 1) and name_base != "image" else f"{name_base}_{suffix}"
        if slot_type == "image" and name == "image":
            name = f"image_{suffix}"

        return {
            "name": name,
            "type": slot_type,
            "data": [sample_text] if sample_text else [],
            "semantic_role": role,
            "source": source,
            "placeholder_type": placeholder_type,
            "placeholder_idx": placeholder_idx,
            "shape_name": raw_name,
            "left_pct": left_pct,
            "top_pct": top_pct,
            "width_pct": width_pct,
            "height_pct": height_pct,
            "position": _shape_position_label(left_pct, top_pct, width_pct, height_pct),
            "font_name": font_name,
            "font_size": font_size,
            "is_content": slot_type != "image" or role == "image",
        }

    def _editable_slots_for_slide(
        self,
        slide: Any,
        layout_capabilities: dict[str, dict[str, Any]],
    ) -> list[dict[str, Any]]:
        layout_name = str(getattr(getattr(slide, "slide_layout", None), "name", "") or "").strip()
        capability = layout_capabilities.get(layout_name)
        if not capability:
            return []
        return [dict(slot) for slot in capability.get("editable_slots", []) or []]

    def _infer_text_kind(self, shape: Any, text: str, shape_idx: int) -> str:
        placeholder = getattr(shape, "placeholder_format", None)
        placeholder_type = str(getattr(placeholder, "type", "") or "").lower()
        if "title" in placeholder_type or shape_idx == 0:
            return "h1"
        if "\n" in text or len(text) > 120:
            return "p"
        return "span"

    def _infer_content_type(self, text_fragments: list[str], image_paths: list[str]) -> str:
        if not text_fragments and image_paths:
            return "image"
        if image_paths and len(image_paths) >= max(1, len(text_fragments)):
            return "mixed-image"
        if len(text_fragments) > 3:
            return "text"
        return "mixed"

    def _write_image_asset(self, image: Any, images_dir: Path, slide_idx: int, shape_idx: int) -> str:
        try:
            blob = getattr(image, "blob", None)
            if not blob:
                return ""
            ext = str(getattr(image, "ext", "") or "png").lstrip(".").lower()
            digest = hashlib.sha1(blob).hexdigest()[:10]
            base_name = str(getattr(image, "filename", "") or f"slide_{slide_idx}_{shape_idx + 1}.{ext}")
            safe_stem = re.sub(r"[^a-zA-Z0-9_.-]+", "_", Path(base_name).stem)[:48] or f"image_{slide_idx}_{shape_idx + 1}"
            file_name = f"{safe_stem}_{digest}.{ext}"
            out_path = images_dir / file_name
            if not out_path.exists():
                out_path.write_bytes(blob)
            return file_name
        except Exception:
            return ""

    def _update_image_stats(
        self,
        *,
        image_stats: dict[str, Any],
        image_name: str,
        slide_idx: int,
        slide_width: float,
        slide_height: float,
        shape_width: float,
        shape_height: float,
    ) -> None:
        entry = image_stats.setdefault(
            image_name,
            {
                "appear_times": 0,
                "relative_area": 0.0,
                "size": [0, 0],
                "slide_numbers": [],
            },
        )
        entry["appear_times"] = int(entry.get("appear_times", 0)) + 1
        area = 0.0
        if slide_width > 0 and slide_height > 0:
            area = (shape_width * shape_height) / (slide_width * slide_height) * 100
        entry["relative_area"] = max(float(entry.get("relative_area", 0.0)), round(area, 2))
        entry["slide_numbers"] = sorted(set(entry.get("slide_numbers", []) + [slide_idx]))
        entry["size"] = [round(shape_width, 2), round(shape_height, 2)]

    def _build_layout_induction(
        self,
        slide_records: list[dict[str, Any]],
        prs: Any,
        layout_capabilities: dict[str, dict[str, Any]],
    ) -> dict[str, Any]:
        if not slide_records:
            return {"functional_keys": [], "language": {"lid": "unknown"}}

        groups: dict[str, dict[str, Any]] = defaultdict(lambda: {"slides": [], "template_id": 0, "elements": []})
        total = len(slide_records)

        for record in slide_records:
            idx = record["index"]
            if idx == 1:
                key = "opening"
            elif idx == total:
                key = "ending"
            elif idx == 2 and total > 4 and self._looks_like_toc(record["text"]):
                key = "table of contents"
            else:
                key = self._content_group_key(record)

            group = groups[key]
            group["slides"].append(idx)
            if not group["template_id"]:
                group["template_id"] = idx
                group["elements"] = record["elements"]
                group["sample_elements"] = record.get("sample_elements", [])
                group["editable_slots"] = record.get("editable_slots", [])
                group["source_layout_name"] = record.get("layout_name", "")
                group["text_count"] = record.get("text_count", 0)
                group["sample_text_count"] = record.get("sample_text_count", 0)
                group["image_count"] = record.get("image_count", 0)

        functional_keys = [key for key in ("opening", "table of contents", "ending") if key in groups]
        layout_induction: dict[str, Any] = {key: value for key, value in groups.items()}

        capability_layouts: dict[str, Any] = {}
        for layout_name, capability in layout_capabilities.items():
            slots = capability.get("editable_slots", []) or []
            if not slots:
                continue
            if not any(slot.get("type") == "text" for slot in slots):
                continue
            capability_key = f"layout:{layout_name}"
            if capability_key in layout_induction:
                continue
            capability_layouts[capability_key] = {
                "slides": [],
                "template_id": 0,
                "elements": slots,
                "editable_slots": slots,
                "source_layout_name": layout_name,
                "source": "layout_placeholder",
                "text_count": capability.get("text_slot_count", 0),
                "image_count": capability.get("image_slot_count", 0),
            }
        layout_induction.update(capability_layouts)
        if capability_layouts:
            layout_induction["layout_capabilities"] = {
                name: {
                    "layout_index": value.get("layout_index"),
                    "text_slot_count": value.get("text_slot_count", 0),
                    "image_slot_count": value.get("image_slot_count", 0),
                    "slot_names": [slot.get("name", "") for slot in value.get("editable_slots", [])],
                }
                for name, value in layout_capabilities.items()
            }
        layout_induction["functional_keys"] = functional_keys
        layout_induction["language"] = self._detect_language_from_records(slide_records)
        return layout_induction

    def _content_group_key(self, record: dict[str, Any]) -> str:
        layout_name = record.get("layout_name") or "content"
        image_count = record.get("editable_image_count", 0) or record.get("image_count", 0)
        text_count = record.get("editable_text_count", 0) or record.get("text_count", 0)
        density = "dense" if text_count >= 4 else "light" if text_count <= 1 else "mid"
        return f"{layout_name}:{density}:{text_count}t:{image_count}i"

    def _build_semantic_model(
        self,
        *,
        slide_records: list[dict[str, Any]],
        layout_induction: dict[str, Any],
        layout_capabilities: dict[str, dict[str, Any]],
        aspect_ratio: str,
    ) -> TemplateSemanticModel:
        sample_elements_by_layout: dict[str, list[dict[str, Any]]] = {}
        decorative_elements: list[dict[str, Any]] = []
        shell_shapes_by_layout: dict[str, list[dict[str, Any]]] = defaultdict(list)
        suspicious: list[str] = []

        for record in slide_records:
            layout_name = str(record.get("layout_name", "") or "")
            if layout_name:
                sample_elements_by_layout.setdefault(layout_name, []).extend(
                    list(record.get("sample_elements", []) or [])
                )
            shape_geo = record.get("shape_geometry", {}) or {}
            shell_candidates: list[dict[str, Any]] = []
            decorative_elements.extend(
                shape
                for shape in shape_geo.get("text_shapes", []) or []
                if not shape.get("is_content")
            )
            shell_candidates.extend(
                shape
                for shape in shape_geo.get("text_shapes", []) or []
                if not shape.get("is_content")
            )
            decorative_elements.extend(
                shape
                for shape in shape_geo.get("image_shapes", []) or []
                if shape.get("semantic_role") in {"background", "decorative"}
            )
            shell_candidates.extend(
                shape
                for shape in shape_geo.get("image_shapes", []) or []
                if shape.get("semantic_role") in {"background", "decorative"}
            )
            if layout_name:
                shell_shapes_by_layout[layout_name].extend(shell_candidates)

        canonical_layouts: list[CanonicalLayout] = []
        layout_capability_payload: dict[str, Any] = {}
        sample_content_payload: dict[str, Any] = {}
        shell_payloads_by_layout: dict[str, dict[str, Any]] = {}
        typography_defaults = self._estimate_typography_defaults(slide_records, layout_induction)

        for key, value in layout_induction.items():
            if key in {"functional_keys", "language", "layout_capabilities"} or not isinstance(value, dict):
                continue
            source_layout_name = str(value.get("source_layout_name", "") or "")
            editable_slots = list(value.get("editable_slots", []) or value.get("elements", []) or [])
            sample_slots_raw = list(value.get("sample_elements", []) or sample_elements_by_layout.get(source_layout_name, []) or [])
            canonical_name = self._canonical_layout_name(
                key=key,
                source_layout_name=source_layout_name,
                editable_slots=editable_slots,
                sample_slots=sample_slots_raw,
            )
            canonical_slots = [
                self._to_canonical_slot(slot, source="layout_capability", typography_defaults=typography_defaults)
                for slot in editable_slots
                if str(slot.get("semantic_role", "") or "").lower() != "decorative"
            ]
            sample_slots = [
                self._to_canonical_slot(slot, source="sample_content", typography_defaults=typography_defaults)
                for slot in sample_slots_raw
                if str(slot.get("semantic_role", "") or "").lower() != "decorative"
            ]

            if not canonical_slots and sample_slots:
                canonical_slots = [
                    slot for slot in sample_slots
                    if slot.type in {"text", "image", "chart", "table"}
                ]
                suspicious.append(
                    f"{canonical_name}: missing editable placeholders, fell back to sample content."
                )

            if not any(slot.type == "text" for slot in canonical_slots):
                suspicious.append(f"{canonical_name}: canonical layout has no editable text slot.")

            archetype = _layout_archetype_from_slots(editable_slots or sample_slots_raw)
            visual_pattern = _visual_pattern_from_slots(editable_slots or sample_slots_raw)
            capability_signature = _capability_signature_from_slots(
                editable_slots or sample_slots_raw
            )
            surface_mode = _surface_mode_from_signature(capability_signature)
            allowed_text_roles = _allowed_text_roles_from_slots(editable_slots or sample_slots_raw)
            allowed_visual_roles = _allowed_visual_roles_from_slots(editable_slots or sample_slots_raw)
            shell_payload = self._build_layout_shell_payload(
                canonical_name=canonical_name,
                source_layout_name=source_layout_name,
                sample_slide_indices=[int(idx) for idx in value.get("slides", []) or []],
                shell_shapes=shell_shapes_by_layout.get(source_layout_name, []),
                editable_slots=editable_slots or sample_slots_raw,
                surface_mode=surface_mode,
            )
            shell_payloads_by_layout[canonical_name] = shell_payload
            layout = CanonicalLayout(
                id=_sanitize_layout_token(canonical_name),
                name=canonical_name,
                aliases=[key] if canonical_name != key else [],
                source_layout_name=source_layout_name,
                layout_archetype=archetype,
                visual_pattern=visual_pattern,
                slot_count=len(canonical_slots),
                slot_roles=[
                    role
                    for role in dict.fromkeys(slot.role for slot in canonical_slots if slot.role)
                ],
                slots=canonical_slots,
                sample_content=sample_slots,
                sample_slide_indices=[int(idx) for idx in value.get("slides", []) or []],
                style_notes=self._layout_style_notes(editable_slots or sample_slots_raw),
                capability_signature=capability_signature,
                surface_mode=surface_mode,
                shell_id=str(shell_payload.get("shell_id", "") or ""),
                shell_mode=str(shell_payload.get("shell_mode", "") or ""),
                protected_regions=list(shell_payload.get("protected_regions", []) or []),
                surface_policy=str(shell_payload.get("surface_policy", "") or ""),
                allowed_text_roles=allowed_text_roles,
                allowed_visual_roles=allowed_visual_roles,
                supports_visual_asset=bool(allowed_visual_roles),
                confidence=0.92 if editable_slots else 0.62,
            )
            canonical_layouts.append(layout)
            layout_capability_payload[canonical_name] = {
                "source_layout_name": source_layout_name,
                "slot_count": len(canonical_slots),
                "slot_roles": layout.slot_roles,
                "layout_archetype": archetype,
                "visual_pattern": visual_pattern,
                "capability_signature": capability_signature,
                "surface_mode": surface_mode,
                "shell_id": shell_payload.get("shell_id", ""),
                "shell_mode": shell_payload.get("shell_mode", ""),
                "surface_policy": shell_payload.get("surface_policy", ""),
                "protected_regions": shell_payload.get("protected_regions", []),
                "allowed_text_roles": allowed_text_roles,
                "allowed_visual_roles": allowed_visual_roles,
                "supports_visual_asset": bool(allowed_visual_roles),
            }
            sample_content_payload[canonical_name] = {
                "sample_slide_indices": layout.sample_slide_indices,
                "sample_values": {
                    slot.name: slot.sample_values[:3]
                    for slot in sample_slots
                    if slot.sample_values
                },
            }

        canonical_layouts = self._merge_duplicate_canonical_layouts(canonical_layouts)
        shell_payloads_by_layout = {
            layout.name: {
                **shell_payloads_by_layout.get(layout.name, {}),
                "shell_id": layout.shell_id,
                "shell_mode": layout.shell_mode,
                "protected_regions": layout.protected_regions,
                "surface_policy": layout.surface_policy,
                "sample_slide_indices": layout.sample_slide_indices,
            }
            for layout in canonical_layouts
        }
        layout_capability_payload = {
            layout.name: {
                "source_layout_name": layout.source_layout_name,
                "slot_count": len(layout.slots),
                "slot_roles": layout.slot_roles,
                "layout_archetype": layout.layout_archetype,
                "visual_pattern": layout.visual_pattern,
                "capability_signature": layout.capability_signature,
                "surface_mode": layout.surface_mode,
                "shell_id": layout.shell_id,
                "shell_mode": layout.shell_mode,
                "surface_policy": layout.surface_policy,
                "protected_regions": layout.protected_regions,
                "allowed_text_roles": layout.allowed_text_roles,
                "allowed_visual_roles": layout.allowed_visual_roles,
                "supports_visual_asset": layout.supports_visual_asset,
            }
            for layout in canonical_layouts
        }
        sample_content_payload = {
            layout.name: {
                "sample_slide_indices": layout.sample_slide_indices,
                "sample_values": {
                    slot.name: slot.sample_values[:3]
                    for slot in layout.sample_content
                    if slot.sample_values
                },
            }
            for layout in canonical_layouts
        }

        decorative_system = self._build_decorative_system(decorative_elements)
        template_shell = self._build_template_shell_model(shell_payloads_by_layout)
        style_system = {
            "aspect_ratio": aspect_ratio,
            "language": layout_induction.get("language", {}),
            "functional_keys": list(layout_induction.get("functional_keys", []) or []),
            "typography_defaults": typography_defaults,
            "layout_surface_modes": {
                layout.name: layout.surface_mode
                for layout in canonical_layouts
            },
            "layout_shell_modes": {
                layout.name: layout.shell_mode
                for layout in canonical_layouts
            },
        }

        return TemplateSemanticModel(
            canonical_layouts=canonical_layouts,
            layout_capability=layout_capability_payload,
            template_shell=template_shell,
            sample_content=sample_content_payload,
            decorative_system=decorative_system,
            style_system=style_system,
            suspicious_ambiguities=sorted(dict.fromkeys(suspicious)),
            low_confidence=any("no editable text slot" in issue for issue in suspicious),
        )

    def _canonical_layout_name(
        self,
        *,
        key: str,
        source_layout_name: str,
        editable_slots: list[dict[str, Any]],
        sample_slots: list[dict[str, Any]],
    ) -> str:
        lowered = key.lower().strip()
        if lowered in {"opening", "ending", "table of contents"}:
            return key
        slot_basis = editable_slots or sample_slots
        archetype = _layout_archetype_from_slots(slot_basis)
        visual_pattern = _visual_pattern_from_slots(slot_basis)
        capability_signature = _capability_signature_from_slots(slot_basis)
        role_counts: Counter[str] = Counter(
            str(slot.get("semantic_role", "") or "slot")
            for slot in slot_basis
            if str(slot.get("semantic_role", "") or "").strip()
        )
        title_part = source_layout_name or key.split(":")[0]
        source_fragment = _sanitize_layout_token(title_part)[:32]
        if not source_fragment or source_fragment in {"blank", "layout"}:
            alias_basis = key.split(":")[0]
            source_fragment = _sanitize_layout_token(alias_basis)[:32]
        role_fragment = "_".join(
            f"{_sanitize_layout_token(role)}{count}"
            for role, count in role_counts.items()
            if role and count > 0
        )[:48]
        fragments = [
            source_fragment,
            _sanitize_layout_token(archetype),
            _sanitize_layout_token(visual_pattern),
        ]
        if capability_signature.get("vertical_title"):
            fragments.append("vertical")
        if capability_signature.get("supports_section_divider"):
            fragments.append("section")
        if capability_signature.get("supports_toc"):
            fragments.append("toc")
        if capability_signature.get("supports_multi_column"):
            fragments.append("multi_column")
        elif capability_signature.get("body_slot_count", 0) >= 2:
            fragments.append("split_body")
        elif capability_signature.get("supports_long_body"):
            fragments.append("long_body")
        if capability_signature.get("supports_figure"):
            fragments.append("figure")
        if capability_signature.get("supports_table"):
            fragments.append("table")
        if capability_signature.get("supports_chart"):
            fragments.append("chart")
        if role_fragment:
            fragments.append(role_fragment)
        normalized = "_".join(part for part in fragments if part)
        return normalized or _sanitize_layout_token(key)

    def _to_canonical_slot(self, slot: dict[str, Any], *, source: str) -> CanonicalSlot:
        return self._to_canonical_slot(slot, source=source, typography_defaults={})

    def _to_canonical_slot(
        self,
        slot: dict[str, Any],
        *,
        source: str,
        typography_defaults: dict[str, dict[str, Any]],
    ) -> CanonicalSlot:
        role = str(slot.get("semantic_role", "") or "")
        font_name = str(slot.get("font_name", "") or "")
        font_size = _as_int(slot.get("font_size", 0), 0)
        fallback = typography_defaults.get(role, {})
        if not font_name:
            font_name = str(fallback.get("font_name", "") or "")
        if font_size <= 0:
            font_size = _as_int(fallback.get("font_size", 0), 0)
        return CanonicalSlot(
            name=str(slot.get("name", "") or ""),
            type=str(slot.get("type", "text") or "text"),
            role=role,
            geometry={
                "left_pct": float(slot.get("left_pct", 0) or 0),
                "top_pct": float(slot.get("top_pct", 0) or 0),
                "width_pct": float(slot.get("width_pct", 0) or 0),
                "height_pct": float(slot.get("height_pct", 0) or 0),
            },
            capacity_hint=_capacity_hint_from_slot({**slot, "font_size": font_size}),
            sample_values=[str(item or "") for item in slot.get("data", []) or [] if str(item or "").strip()],
            font_name=font_name,
            font_size=font_size,
            source=source,
            placeholder_type=str(slot.get("placeholder_type", "") or ""),
            placeholder_idx=slot.get("placeholder_idx"),
            shape_name=str(slot.get("shape_name", "") or ""),
            notes=[str(slot.get("position", "") or "")] if slot.get("position") else [],
        )

    def _estimate_typography_defaults(
        self,
        slide_records: list[dict[str, Any]],
        layout_induction: dict[str, Any],
    ) -> dict[str, dict[str, Any]]:
        role_fonts: dict[str, list[str]] = defaultdict(list)
        role_sizes: dict[str, list[int]] = defaultdict(list)
        for record in slide_records:
            geo = record.get("shape_geometry", {}) or {}
            for shape in geo.get("text_shapes", []) or []:
                role = str(shape.get("semantic_role", "") or "")
                if not role or role == "decorative":
                    continue
                font_name = str(shape.get("font_name", "") or "")
                font_size = _as_int(shape.get("font_size", 0), 0)
                if font_name:
                    role_fonts[role].append(font_name)
                if font_size > 0:
                    role_sizes[role].append(font_size)
        defaults: dict[str, dict[str, Any]] = {}
        for role in set(role_fonts.keys()) | set(role_sizes.keys()):
            entry: dict[str, Any] = {}
            if role_fonts.get(role):
                entry["font_name"] = Counter(role_fonts[role]).most_common(1)[0][0]
            if role_sizes.get(role):
                sizes = sorted(role_sizes[role])
                entry["font_size"] = sizes[len(sizes) // 2]
            defaults[role] = entry
        if "body" not in defaults:
            defaults["body"] = {"font_size": 18}
        if "title" not in defaults:
            defaults["title"] = {"font_size": 28}
        return defaults

    def _merge_duplicate_canonical_layouts(
        self,
        layouts: list[CanonicalLayout],
    ) -> list[CanonicalLayout]:
        merged: dict[tuple, CanonicalLayout] = {}
        order: list[tuple] = []
        for layout in layouts:
            signature = (
                layout.name,
                layout.layout_archetype,
                layout.visual_pattern,
                json.dumps(layout.capability_signature, ensure_ascii=False, sort_keys=True),
                _layout_signature_from_slots(layout.slots),
            )
            if signature not in merged:
                merged[signature] = layout
                order.append(signature)
                continue
            current = merged[signature]
            current.aliases = list(dict.fromkeys([*current.aliases, *layout.aliases, layout.source_layout_name]).keys())
            current.sample_slide_indices = sorted(set(current.sample_slide_indices + layout.sample_slide_indices))
            current.style_notes = list(dict.fromkeys([*current.style_notes, *layout.style_notes]).keys())
            current.confidence = max(current.confidence, layout.confidence)
            for slot in layout.sample_content:
                if slot.name not in {existing.name for existing in current.sample_content}:
                    current.sample_content.append(slot)
        merged_list = [merged[key] for key in order]
        return self._ensure_unique_layout_names(merged_list)

    def _ensure_unique_layout_names(
        self,
        layouts: list[CanonicalLayout],
    ) -> list[CanonicalLayout]:
        grouped: dict[str, list[CanonicalLayout]] = defaultdict(list)
        for layout in layouts:
            grouped[layout.name].append(layout)

        result: list[CanonicalLayout] = []
        for name, items in grouped.items():
            if len(items) == 1:
                result.extend(items)
                continue
            for index, layout in enumerate(
                sorted(
                    items,
                    key=lambda item: (
                        len(item.slots),
                        len(item.sample_slide_indices),
                        item.source_layout_name,
                    ),
                ),
                start=1,
            ):
                layout.name = f"{name}_v{index}"
                result.append(layout)
        return result

    def _layout_style_notes(self, slots: list[dict[str, Any]]) -> list[str]:
        notes: list[str] = []
        title_slots = [slot for slot in slots if slot.get("semantic_role") == "title"]
        body_slots = [slot for slot in slots if slot.get("semantic_role") == "body"]
        if title_slots:
            max_title = max(_as_int(slot.get("font_size", 0), 0) for slot in title_slots)
            if max_title:
                notes.append(f"title_font_size≈{max_title}px")
        if body_slots:
            body_sizes = [_as_int(slot.get("font_size", 0), 0) for slot in body_slots if _as_int(slot.get("font_size", 0), 0) > 0]
            if body_sizes:
                body_sizes.sort()
                notes.append(f"body_font_size≈{body_sizes[len(body_sizes)//2]}px")
        return notes

    def _build_layout_shell_payload(
        self,
        *,
        canonical_name: str,
        source_layout_name: str,
        sample_slide_indices: list[int],
        shell_shapes: list[dict[str, Any]],
        editable_slots: list[dict[str, Any]],
        surface_mode: str,
    ) -> dict[str, Any]:
        backgrounds = [
            shape for shape in shell_shapes
            if shape.get("type") == "image" and shape.get("semantic_role") == "background"
        ]
        decorative_shapes = [
            shape for shape in shell_shapes
            if shape.get("semantic_role") in {"decorative", "number_badge", "section_label"}
            or (shape.get("type") == "text" and not shape.get("is_content"))
        ]
        protected_regions = [
            region
            for region in (
                _protected_region_from_decorative_shape(shape)
                for shape in decorative_shapes
            )
            if region is not None
        ]
        title_slots = [
            slot for slot in editable_slots
            if str(slot.get("semantic_role", "") or "") in {"title", "subtitle", "section_label"}
        ]
        for slot in title_slots[:3]:
            region = _region_from_shape(slot, role="title_slot")
            if region not in protected_regions and region.get("top_pct", 100) <= 25:
                protected_regions.append(region)

        background_assets = [
            {
                "asset_name": str(shape.get("asset_name", "") or ""),
                "geometry": _region_from_shape(shape, role="background"),
                "position": str(shape.get("position", "") or ""),
                "shape_path": str(shape.get("shape_path", "") or ""),
            }
            for shape in backgrounds
            if str(shape.get("asset_name", "") or "").strip()
        ]
        decorative_layers = [
            _shape_shell_layer(shape, layer_type=str(shape.get("type", "") or "shape"), z_index=index)
            for index, shape in enumerate(decorative_shapes[:24], start=1)
        ]

        if background_assets:
            shell_mode = NATIVE_BACKGROUND
        elif decorative_layers:
            shell_mode = LAYERED_SHELL
        else:
            shell_mode = NO_SHELL

        if shell_mode == LAYERED_SHELL and len(decorative_layers) >= 18:
            shell_mode = RASTER_SHELL

        has_body = any(str(slot.get("semantic_role", "") or "") == "body" for slot in editable_slots)
        has_visual = any(str(slot.get("type", "") or "") in {"image", "chart", "table"} for slot in editable_slots)
        if surface_mode in {"figure_panel", "light_surface"} or has_body or has_visual:
            surface_policy = "adaptive_surface_allowed"
        else:
            surface_policy = "shell_text_only"

        notes: list[str] = []
        if background_assets:
            notes.append("full-slide image background can be replayed as PPTX background")
        if decorative_layers:
            notes.append("decorative layers/protected regions should be preserved before content overlay")
        if surface_policy == "adaptive_surface_allowed":
            notes.append("body, formula, table, and figure evidence may use a readable content surface")
        if shell_mode == RASTER_SHELL:
            notes.append("complex decorative shell should be rasterized when a renderer is available")

        return {
            "shell_id": f"shell_{_sanitize_layout_token(canonical_name)}",
            "shell_mode": shell_mode,
            "source_layout_name": source_layout_name,
            "sample_slide_indices": sample_slide_indices,
            "background_assets": background_assets,
            "decorative_layers": decorative_layers,
            "protected_regions": protected_regions[:16],
            "safe_content_regions": self._safe_content_regions_from_slots(editable_slots),
            "surface_policy": surface_policy,
            "reuse_confidence": 0.9 if shell_mode in {NATIVE_BACKGROUND, LAYERED_SHELL} else 0.55 if shell_mode == RASTER_SHELL else 0.2,
            "notes": notes,
        }

    def _safe_content_regions_from_slots(self, slots: list[dict[str, Any]]) -> list[dict[str, Any]]:
        regions: list[dict[str, Any]] = []
        for slot in slots:
            role = str(slot.get("semantic_role", "") or "")
            slot_type = str(slot.get("type", "") or "")
            if role in {"body", "caption"} or slot_type in {"image", "chart", "table"}:
                regions.append(_region_from_shape(slot, role=role or slot_type or "content"))
        return regions[:8]

    def _build_template_shell_model(
        self,
        shell_payloads_by_layout: dict[str, dict[str, Any]],
    ) -> dict[str, Any]:
        mode_counts = Counter(
            str(payload.get("shell_mode", "") or NO_SHELL)
            for payload in shell_payloads_by_layout.values()
        )
        return {
            "version": 1,
            "summary": {
                "layout_count": len(shell_payloads_by_layout),
                "shell_mode_counts": dict(mode_counts),
                "has_replayable_shell": any(
                    mode in {NATIVE_BACKGROUND, LAYERED_SHELL, RASTER_SHELL}
                    for mode in mode_counts
                ),
            },
            "layouts": shell_payloads_by_layout,
        }

    def _build_decorative_system(self, decorative_elements: list[dict[str, Any]]) -> dict[str, Any]:
        background_count = sum(1 for item in decorative_elements if item.get("semantic_role") == "background")
        decorative_count = sum(1 for item in decorative_elements if item.get("semantic_role") == "decorative")
        repeated_positions = Counter(
            str(item.get("position", "") or "")
            for item in decorative_elements
            if str(item.get("position", "") or "").strip()
        )
        return {
            "background_count": background_count,
            "decorative_count": decorative_count,
            "repeated_positions": [
                {"position": position, "count": count}
                for position, count in repeated_positions.most_common(8)
            ],
        }

    def _looks_like_toc(self, text: str) -> bool:
        lowered = text.lower()
        return any(term in lowered for term in ("目录", "contents", "agenda", "outline", "table of contents"))

    def _detect_language_from_records(self, slide_records: list[dict[str, Any]]) -> dict[str, str]:
        texts = [record.get("text", "") for record in slide_records if record.get("text")]
        return _detect_language(texts)

    def _most_common_fonts(self, font_names: list[str]) -> list[str]:
        if not font_names:
            return []
        counter = Counter(font_names)
        return [name for name, _ in counter.most_common(6)]

    def _get_aspect_ratio(self, slide_width: float, slide_height: float) -> str:
        if not slide_width or not slide_height:
            return "16:9"
        ratio = slide_width / slide_height
        if abs(ratio - 16 / 9) < 0.15:
            return "16:9"
        if abs(ratio - 4 / 3) < 0.15:
            return "4:3"
        return f"{round(slide_width)}:{round(slide_height)}"

    def _persist_analysis_outputs(self, output_dir: Path, analysis: TemplateAnalysis) -> None:
        try:
            (output_dir / "analysis.json").write_text(
                json.dumps(analysis.to_dict(), ensure_ascii=False, indent=2),
                encoding="utf-8",
            )
        except Exception as exc:
            logger.debug("Failed to write analysis.json: %s", exc)
        try:
            (output_dir / "slide_induction.json").write_text(
                json.dumps(analysis.layout_induction, ensure_ascii=False, indent=2),
                encoding="utf-8",
            )
        except Exception as exc:
            logger.debug("Failed to write slide_induction.json: %s", exc)
        try:
            (output_dir / "semantic_model.json").write_text(
                json.dumps(analysis.semantic_model, ensure_ascii=False, indent=2),
                encoding="utf-8",
            )
        except Exception as exc:
            logger.debug("Failed to write semantic_model.json: %s", exc)
        try:
            (output_dir / "template_audit.json").write_text(
                json.dumps(_template_audit_payload(analysis.semantic_model), ensure_ascii=False, indent=2),
                encoding="utf-8",
            )
        except Exception as exc:
            logger.debug("Failed to write template_audit.json: %s", exc)
        try:
            template_shell = {}
            if isinstance(analysis.semantic_model, dict):
                template_shell = analysis.semantic_model.get("template_shell", {}) or {}
            (output_dir / "template_shell.json").write_text(
                json.dumps(template_shell, ensure_ascii=False, indent=2),
                encoding="utf-8",
            )
        except Exception as exc:
            logger.debug("Failed to write template_shell.json: %s", exc)
        try:
            template_shell = {}
            if isinstance(analysis.semantic_model, dict):
                template_shell = analysis.semantic_model.get("template_shell", {}) or {}
            audit = {
                "summary": template_shell.get("summary", {}) if isinstance(template_shell, dict) else {},
                "layout_checks": [
                    {
                        "layout": name,
                        "shell_id": payload.get("shell_id", ""),
                        "shell_mode": payload.get("shell_mode", ""),
                        "has_background": bool(payload.get("background_assets")),
                        "protected_region_count": len(payload.get("protected_regions", []) or []),
                        "surface_policy": payload.get("surface_policy", ""),
                        "safe_content_region_count": len(payload.get("safe_content_regions", []) or []),
                    }
                    for name, payload in ((template_shell.get("layouts", {}) or {}).items() if isinstance(template_shell, dict) else [])
                    if isinstance(payload, dict)
                ],
            }
            (output_dir / "template_shell_audit.json").write_text(
                json.dumps(audit, ensure_ascii=False, indent=2),
                encoding="utf-8",
            )
        except Exception as exc:
            logger.debug("Failed to write template_shell_audit.json: %s", exc)
        try:
            (output_dir / "image_stats.json").write_text(
                json.dumps(analysis.image_stats, ensure_ascii=False, indent=2),
                encoding="utf-8",
            )
        except Exception as exc:
            logger.debug("Failed to write image_stats.json: %s", exc)
        try:
            (output_dir / "shape_geometry.json").write_text(
                json.dumps(analysis.shape_geometry, ensure_ascii=False, indent=2),
                encoding="utf-8",
            )
        except Exception as exc:
            logger.debug("Failed to write shape_geometry.json: %s", exc)
        try:
            (output_dir / "description.txt").write_text(
                f"{analysis.template_name}\nAspect ratio: {analysis.aspect_ratio}, {analysis.slide_count} slides.",
                encoding="utf-8",
            )
        except Exception as exc:
            logger.debug("Failed to write description.txt: %s", exc)

    def cleanup(self, pptx_path: Path | str):
        pptx_path = Path(pptx_path)
        temp_dir = self.workspace / ".template_analysis" / pptx_path.stem
        if temp_dir.exists():
            shutil.rmtree(temp_dir, ignore_errors=True)


__all__ = ["TemplateAnalyzer", "TemplateAnalysis"]

"""
TemplateExtractors — 模板提取器合集

合并自:
- StyleExtractor (颜色/字体/间距)
- ContentPatternExtractor (标题/正文/图表模式)

两个提取器共享同一输入 (slides_html)，在 main.py 中总是一起调用。
"""

from __future__ import annotations

import colorsys
import logging
import re
from collections import Counter
from typing import Any, Optional

from ..core.template_models import (
    ColorPalette,
    ContentPatterns,
    DesignConstraints,
    Typography,
)

logger = logging.getLogger(__name__)


def _hex_to_rgb(color: str) -> tuple[int, int, int] | None:
    color = str(color or "").strip().lower()
    if color.startswith("#") and len(color) == 4:
        color = f"#{color[1]*2}{color[2]*2}{color[3]*2}"
    if color.startswith("#") and len(color) == 7:
        try:
            return int(color[1:3], 16), int(color[3:5], 16), int(color[5:7], 16)
        except Exception:
            return None
    return None


def _relative_luminance(color: str) -> float | None:
    rgb = _hex_to_rgb(color)
    if rgb is None:
        return None

    def _channel(v: int) -> float:
        x = v / 255.0
        return x / 12.92 if x <= 0.04045 else ((x + 0.055) / 1.055) ** 2.4

    r, g, b = (_channel(v) for v in rgb)
    return 0.2126 * r + 0.7152 * g + 0.0722 * b


def _contrast_ratio(c1: str, c2: str) -> float:
    l1 = _relative_luminance(c1)
    l2 = _relative_luminance(c2)
    if l1 is None or l2 is None:
        return 0.0
    lighter = max(l1, l2)
    darker = min(l1, l2)
    return (lighter + 0.05) / (darker + 0.05)


def _is_light_color(color: str) -> bool:
    lum = _relative_luminance(color)
    return bool(lum is not None and lum >= 0.55)


def _pick_best_contrast(base_color: str, candidates: list[str]) -> str:
    ranked = sorted(
        (
            (candidate, _contrast_ratio(base_color, candidate))
            for candidate in candidates
            if candidate
        ),
        key=lambda item: item[1],
        reverse=True,
    )
    return ranked[0][0] if ranked else ""


def _pick_same_tone_fallback(base_color: str, *, prefer_light: bool) -> str:
    rgb = _hex_to_rgb(base_color)
    if rgb is None:
        return "#ffffff" if prefer_light else "#1a1a1a"
    r, g, b = rgb
    hue, sat, val = colorsys.rgb_to_hsv(r / 255, g / 255, b / 255)
    target_val = 0.96 if prefer_light else 0.14
    target_sat = min(0.22 if prefer_light else sat, sat)
    rr, gg, bb = colorsys.hsv_to_rgb(hue, target_sat, target_val)
    return f"#{round(rr*255):02x}{round(gg*255):02x}{round(bb*255):02x}"


def _iter_shape_tree(shapes: Any):
    for shape in list(shapes or []):
        yield shape
        child_shapes = getattr(shape, "shapes", None)
        if child_shapes is not None:
            yield from _iter_shape_tree(child_shapes)


def _placeholder_type(shape: Any) -> str:
    try:
        return str(shape.placeholder_format.type)
    except Exception:
        return ""


def _is_ignored_placeholder(shape: Any) -> bool:
    placeholder_type = _placeholder_type(shape).upper()
    return any(
        token in placeholder_type
        for token in ("DATE", "FOOTER", "SLIDE_NUMBER", "HEADER")
    )


def _iter_text_shapes_from_presentation(prs: Any):
    for slide in getattr(prs, "slides", []) or []:
        for shape in _iter_shape_tree(getattr(slide, "shapes", []) or []):
            if getattr(shape, "has_text_frame", False):
                yield shape, "slide"
    for layout in getattr(prs, "slide_layouts", []) or []:
        for shape in _iter_shape_tree(getattr(layout, "shapes", []) or []):
            if getattr(shape, "has_text_frame", False):
                yield shape, "layout"
    for master in getattr(prs, "slide_masters", []) or []:
        for shape in _iter_shape_tree(getattr(master, "shapes", []) or []):
            if getattr(shape, "has_text_frame", False):
                yield shape, "master"


def _is_title_shape(shape: Any, *, source: str, size_pt: float | None) -> bool:
    placeholder_type = _placeholder_type(shape).upper()
    if "TITLE" in placeholder_type or "CENTER_TITLE" in placeholder_type:
        return True
    try:
        if int(shape.placeholder_format.idx) == 0:
            return True
    except Exception:
        pass
    text = (getattr(shape, "text", "") or "").strip()
    if size_pt is not None and size_pt >= 28 and len(text) <= 120:
        return True
    return False


def _color_to_hex(color: Any) -> str:
    try:
        rgb = getattr(color, "rgb", None)
        if rgb is not None:
            return f"#{int(rgb[0]):02x}{int(rgb[1]):02x}{int(rgb[2]):02x}"
    except Exception:
        pass
    return ""


# ═══════════════════════════════════════════════
# StyleExtractor
# ═══════════════════════════════════════════════


class StyleExtractor:
    """从 PPTX 直接提取设计约束（颜色、字体、间距）

    优先从 PPTX 原始数据提取（稳定、精确），HTML 作为补充来源。
    """

    async def extract(
        self,
        slides_html: list[str],
        pptx_path: Optional[str] = None,
    ) -> DesignConstraints:
        constraints = DesignConstraints()

        # ── 1. 颜色：HTML + PPTX 联合提取 ──
        colors = self._extract_colors_from_html(slides_html)
        if colors:
            constraints.color_palette = colors

        if pptx_path:
            pptx_colors = self._extract_colors_from_pptx(pptx_path)
            for field_name in (
                "background",
                "surface",
                "text",
                "primary_text",
                "inverse_text",
                "muted_text",
                "border",
                "primary",
                "secondary",
                "accent",
            ):
                value = getattr(pptx_colors, field_name, "")
                if value:
                    setattr(constraints.color_palette, field_name, value)
            if pptx_colors.additional:
                existing = {
                    constraints.color_palette.primary,
                    constraints.color_palette.secondary,
                    constraints.color_palette.accent,
                    constraints.color_palette.background,
                    constraints.color_palette.surface,
                    constraints.color_palette.text,
                    constraints.color_palette.primary_text,
                    constraints.color_palette.inverse_text,
                    constraints.color_palette.muted_text,
                    constraints.color_palette.border,
                    *constraints.color_palette.additional,
                } - {""}
                constraints.color_palette.additional.extend(
                    color for color in pptx_colors.additional if color not in existing
                )
            bg_color = await self._extract_background_color(pptx_path)
            if bg_color:
                constraints.color_palette.background = bg_color
            constraints.color_palette = self._finalize_palette_roles(constraints.color_palette)

        # ── 2. 字体 & 字号：优先 PPTX，HTML 补充 ──
        if pptx_path:
            typography = self._extract_typography_from_pptx(pptx_path)
        else:
            typography = Typography()
        # HTML 补充：如果 PPTX 没有提取到某些字段
        html_typography = self._extract_fonts_from_html(slides_html)
        if not typography.title_font and html_typography.title_font:
            typography.title_font = html_typography.title_font
        if not typography.body_font and html_typography.body_font:
            typography.body_font = html_typography.body_font
        if not typography.title_size and html_typography.title_size:
            typography.title_size = html_typography.title_size
        if not typography.body_size and html_typography.body_size:
            typography.body_size = html_typography.body_size
        constraints.typography = typography

        # ── 3. 间距：优先 PPTX，HTML 补充 ──
        if pptx_path:
            spacing = self._extract_spacing_from_pptx(pptx_path)
        else:
            spacing = {}
        if not spacing:
            spacing = self._extract_spacing_from_html(slides_html)
        if spacing:
            constraints.spacing = spacing

        return constraints

    # ── PPTX 直接提取方法 ──

    def _extract_typography_from_pptx(self, pptx_path: str) -> Typography:
        """从 PPTX 直接读取字体和字号（稳定、精确）"""
        typography = Typography()
        try:
            from pptx import Presentation
            from pptx.util import Pt, Emu
            prs = Presentation(str(pptx_path))

            title_fonts: list[str] = []
            body_fonts: list[str] = []
            title_sizes: list[float] = []  # in pt
            body_sizes: list[float] = []   # in pt

            for shape, source in _iter_text_shapes_from_presentation(prs):
                if _is_ignored_placeholder(shape):
                    continue
                placeholder_type = _placeholder_type(shape).upper()
                is_body_placeholder = any(token in placeholder_type for token in ("BODY", "OBJECT", "SUBTITLE"))
                text = (getattr(shape, "text", "") or "").strip()

                # Empty placeholders still carry useful default fonts/styles on
                # layouts/masters.  They do not have run text, so we use their
                # paragraph runs when present and otherwise keep scanning.
                saw_run = False
                try:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if not run.text.strip() and source == "slide":
                                continue
                            saw_run = True
                            # 字号（Emu → pt）
                            font_size = run.font.size
                            size_pt = font_size / Pt(1) if font_size else None
                            effective_title = _is_title_shape(shape, source=source, size_pt=size_pt)

                            # 字体名
                            font_name = run.font.name
                            if font_name:
                                if effective_title:
                                    title_fonts.append(font_name)
                                elif is_body_placeholder or text:
                                    body_fonts.append(font_name)
                                else:
                                    body_fonts.append(font_name)

                            if size_pt is not None:
                                if effective_title:
                                    title_sizes.append(size_pt)
                                elif is_body_placeholder or text:
                                    body_sizes.append(size_pt)
                                else:
                                    body_sizes.append(size_pt)
                except Exception:
                    continue

                # Some placeholders expose font defaults at paragraph level but
                # no populated run.  Treat missing body size as unknown rather
                # than forcing 0; HTML and master defaults may fill it later.
                if not saw_run:
                    continue

            # 如果占位符分类没有区分出标题/正文，用字号阈值分组
            if not title_sizes and not body_sizes:
                all_sizes = title_sizes + body_sizes
                if all_sizes:
                    median = sorted(all_sizes)[len(all_sizes) // 2]
                    title_sizes = [s for s in all_sizes if s >= median]
                    body_sizes = [s for s in all_sizes if s < median]

            if title_fonts:
                typography.title_font = Counter(title_fonts).most_common(1)[0][0]
            if body_fonts:
                typography.body_font = Counter(body_fonts).most_common(1)[0][0]
            if title_sizes:
                title_sizes.sort()
                typography.title_size = int(title_sizes[len(title_sizes) // 2])
            if body_sizes:
                body_sizes.sort()
                typography.body_size = int(body_sizes[len(body_sizes) // 2])
            elif title_sizes:
                smallest_title = min(title_sizes)
                if smallest_title > 18:
                    typography.body_size = max(14, int(round(smallest_title * 0.55)))

            logger.info(
                f"[TEMPLATE] PPTX typography: title_font={typography.title_font!r}, "
                f"title_size={typography.title_size}, body_font={typography.body_font!r}, "
                f"body_size={typography.body_size}"
            )
        except Exception as e:
            logger.warning(f"[TEMPLATE] PPTX typography extraction failed: {e}")

        return typography

    def _extract_spacing_from_pptx(self, pptx_path: str) -> dict:
        """从 PPTX 直接读取间距信息"""
        spacing = {}
        try:
            from pptx import Presentation
            from pptx.util import Emu
            prs = Presentation(str(pptx_path))

            margins: list[float] = []  # in pt
            paddings: list[float] = []  # in pt (internal margin of text frames)

            for shape, _source in _iter_text_shapes_from_presentation(prs):
                if _is_ignored_placeholder(shape):
                    continue
                tf = shape.text_frame
                # 内边距（text frame margin）
                for attr in ('margin_left', 'margin_right', 'margin_top', 'margin_bottom'):
                    val = getattr(tf, attr, None)
                    if val is not None and val > 0:
                        paddings.append(round(val / 12700, 1))  # Emu to pt

                # 段落间距
                for para in tf.paragraphs:
                    pf = getattr(para, 'paragraph_format', None)
                    if pf is None:
                        continue
                    sp_before = getattr(pf, 'space_before', None)
                    sp_after = getattr(pf, 'space_after', None)
                    if sp_before is not None and sp_before > 0:
                        margins.append(round(sp_before / 12700, 1))
                    if sp_after is not None and sp_after > 0:
                        margins.append(round(sp_after / 12700, 1))

            if margins:
                margins.sort()
                spacing['margin'] = f"{margins[len(margins) // 2]}pt"
            if paddings:
                paddings.sort()
                spacing['padding'] = f"{paddings[len(paddings) // 2]}pt"

            if spacing:
                logger.info(f"[TEMPLATE] PPTX spacing: {spacing}")
        except Exception as e:
            logger.warning(f"[TEMPLATE] PPTX spacing extraction failed: {e}")

        return spacing

    # ── 颜色提取（HTML + PPTX 混合）──

    # 颜色归一化：统一为 #rrggbb 小写
    @staticmethod
    def _normalize_color(raw: str) -> str:
        raw = raw.strip()
        if raw.startswith("rgb"):
            m = re.match(r'rgba?\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)', raw)
            if m:
                return f"#{int(m.group(1)):02x}{int(m.group(2)):02x}{int(m.group(3)):02x}"
            return raw.lower()
        raw = raw.lower()
        if raw.startswith("#") and len(raw) == 4:
            return f"#{raw[1]*2}{raw[2]*2}{raw[3]*2}"
        return raw

    # 不参与语义分类的安全/中性颜色
    _NEUTRAL_COLORS = {
        "#ffffff", "#000000", "#fff", "#000",
        "transparent", "inherit", "initial", "currentcolor",
    }

    def _finalize_palette_roles(self, palette: ColorPalette) -> ColorPalette:
        all_candidates = [
            palette.primary,
            palette.secondary,
            palette.accent,
            palette.text,
            palette.background,
            palette.surface,
            palette.primary_text,
            palette.inverse_text,
            palette.muted_text,
            palette.border,
            *(palette.additional or []),
        ]
        unique_candidates = [c for c in dict.fromkeys(c for c in all_candidates if c)]

        if palette.background and not palette.surface:
            palette.surface = _pick_same_tone_fallback(palette.background, prefer_light=True)
        if palette.background and not palette.inverse_text:
            palette.inverse_text = _pick_best_contrast(
                palette.background,
                [palette.text, palette.primary, palette.secondary, palette.accent, *unique_candidates, "#ffffff", "#f9f9f9"],
            ) or "#ffffff"
        if palette.surface and not palette.primary_text:
            palette.primary_text = _pick_best_contrast(
                palette.surface,
                [palette.text, palette.primary, palette.secondary, palette.accent, *unique_candidates, "#111111", "#1a1a1a"],
            ) or "#111111"
        if not palette.text:
            palette.text = palette.primary_text or palette.inverse_text
        if not palette.muted_text:
            palette.muted_text = palette.secondary or palette.accent or palette.text
        if not palette.border:
            palette.border = palette.accent or palette.secondary or palette.background
        reserved = {
            palette.primary,
            palette.secondary,
            palette.accent,
            palette.background,
            palette.text,
            palette.surface,
            palette.primary_text,
            palette.inverse_text,
            palette.muted_text,
            palette.border,
        } - {""}
        palette.additional = [color for color in (palette.additional or []) if color not in reserved]
        return palette

    def _extract_colors_from_html(self, slides_html: list[str]) -> ColorPalette:
        """按 CSS 属性语义分桶提取颜色，再按频率分配角色。

        桶:
          background — background-color / background 属性中的颜色
          text       — color 属性（非 background 上下文）中的颜色
          other      — border / 其他属性中的颜色
        """
        palette = ColorPalette()

        # 正则：捕获 (CSS属性名, 颜色值)
        # 匹配 inline style 和 <style> 块中的颜色声明
        prop_color_re = re.compile(
            r'(background(?:-color)?|(?<!-)color|border(?:-color)?)'
            r'\s*:\s*'
            r'(#[0-9A-Fa-f]{3,6}|rgba?\([^)]+\))',
            re.IGNORECASE,
        )

        bg_colors: list[str] = []
        text_colors: list[str] = []
        other_colors: list[str] = []

        for html in slides_html:
            for prop_match in prop_color_re.finditer(html):
                prop_name = prop_match.group(1).lower()
                raw_color = prop_match.group(2)
                norm = self._normalize_color(raw_color)
                if norm in self._NEUTRAL_COLORS:
                    continue
                if "background" in prop_name:
                    bg_colors.append(norm)
                elif prop_name == "color":
                    text_colors.append(norm)
                else:
                    other_colors.append(norm)

        surface_colors: list[str] = []
        border_colors: list[str] = []

        if bg_colors:
            palette.background = Counter(bg_colors).most_common(1)[0][0]
        if text_colors:
            palette.text = Counter(text_colors).most_common(1)[0][0]
        for color, count in Counter(bg_colors).most_common(6):
            if color != palette.background and _is_light_color(color):
                surface_colors.extend([color] * count)
        for color, count in Counter(other_colors).most_common(6):
            border_colors.extend([color] * count)
        if surface_colors:
            palette.surface = Counter(surface_colors).most_common(1)[0][0]
        if border_colors:
            palette.border = Counter(border_colors).most_common(1)[0][0]

        # 合并所有非中性颜色，排除已分配的 background/text，按频率分配 primary/secondary/accent
        all_semantic = bg_colors + text_colors + other_colors
        assigned = {palette.background, palette.text, palette.surface, palette.border} - {""}
        remaining = [c for c in all_semantic if c not in assigned]
        if remaining:
            ranked = Counter(remaining).most_common(8)
            if ranked:
                palette.primary = ranked[0][0]
            if len(ranked) > 1:
                palette.secondary = ranked[1][0]
            if len(ranked) > 2:
                palette.accent = ranked[2][0]
            # additional: top-4 ~ top-8
            if len(ranked) > 3:
                palette.additional = [c for c, _ in ranked[3:]]

        return self._finalize_palette_roles(palette)

    def _extract_colors_from_pptx(self, pptx_path: str) -> ColorPalette:
        palette = ColorPalette()
        try:
            from pptx import Presentation
            prs = Presentation(str(pptx_path))

            fill_colors: list[str] = []
            surface_fill_colors: list[str] = []
            text_colors: list[str] = []
            body_text_colors: list[str] = []
            title_text_colors: list[str] = []
            light_body_text_colors: list[str] = []
            dark_body_text_colors: list[str] = []
            border_colors: list[str] = []
            large_fill_colors: list[str] = []
            slide_area = float(getattr(prs, "slide_width", 0) or 1) * float(getattr(prs, "slide_height", 0) or 1)

            def add_fill(shape: Any) -> None:
                try:
                    fill = getattr(shape, "fill", None)
                    if fill is not None:
                        color = _color_to_hex(getattr(fill, "fore_color", None))
                        if color and color not in self._NEUTRAL_COLORS:
                            fill_colors.append(color)
                            width = float(getattr(shape, "width", 0) or 0)
                            height = float(getattr(shape, "height", 0) or 0)
                            area_ratio = (width * height / slide_area) if slide_area > 0 else 0
                            if slide_area > 0 and area_ratio >= 0.35:
                                large_fill_colors.append(color)
                            elif area_ratio >= 0.08 and _is_light_color(color):
                                surface_fill_colors.append(color)
                except Exception:
                    pass
                try:
                    line = getattr(shape, "line", None)
                    color = _color_to_hex(getattr(line, "color", None))
                    if color and color not in self._NEUTRAL_COLORS:
                        border_colors.append(color)
                except Exception:
                    pass

            def add_text(shape: Any) -> None:
                tf = getattr(shape, "text_frame", None)
                if tf is None:
                    return
                try:
                    for para in tf.paragraphs:
                        for run in para.runs:
                            if not run.text.strip():
                                continue
                            color = _color_to_hex(getattr(run.font, "color", None))
                            if color and color not in self._NEUTRAL_COLORS:
                                text_colors.append(color)
                                try:
                                    size_pt = run.font.size.pt if run.font.size else None
                                except Exception:
                                    size_pt = None
                                if _is_title_shape(shape, source="slide", size_pt=size_pt):
                                    title_text_colors.append(color)
                                else:
                                    body_text_colors.append(color)
                                    if _is_light_color(color):
                                        light_body_text_colors.append(color)
                                    else:
                                        dark_body_text_colors.append(color)
                except Exception:
                    pass

            for slide in getattr(prs, "slides", []) or []:
                for shape in _iter_shape_tree(getattr(slide, "shapes", []) or []):
                    add_fill(shape)
                    add_text(shape)
            for layout in getattr(prs, "slide_layouts", []) or []:
                for shape in _iter_shape_tree(getattr(layout, "shapes", []) or []):
                    add_fill(shape)
                    add_text(shape)
            for master in getattr(prs, "slide_masters", []) or []:
                for shape in _iter_shape_tree(getattr(master, "shapes", []) or []):
                    add_fill(shape)
                    add_text(shape)

            if large_fill_colors:
                palette.background = Counter(large_fill_colors).most_common(1)[0][0]
            elif fill_colors:
                palette.background = Counter(fill_colors).most_common(1)[0][0]
            if surface_fill_colors:
                palette.surface = Counter(surface_fill_colors).most_common(1)[0][0]
            if dark_body_text_colors:
                palette.primary_text = Counter(dark_body_text_colors).most_common(1)[0][0]
            elif body_text_colors:
                palette.primary_text = Counter(body_text_colors).most_common(1)[0][0]
            if light_body_text_colors:
                palette.inverse_text = Counter(light_body_text_colors).most_common(1)[0][0]
            if body_text_colors:
                palette.text = Counter(body_text_colors).most_common(1)[0][0]
            elif text_colors:
                palette.text = Counter(text_colors).most_common(1)[0][0]
            if border_colors:
                palette.border = Counter(border_colors).most_common(1)[0][0]

            assigned = {
                palette.background,
                palette.surface,
                palette.text,
                palette.primary_text,
                palette.inverse_text,
                palette.border,
            } - {""}
            remaining = [
                color for color in title_text_colors + fill_colors + border_colors + text_colors
                if color not in assigned
            ]
            ranked = Counter(remaining).most_common(8)
            if ranked:
                palette.primary = ranked[0][0]
            if len(ranked) > 1:
                palette.secondary = ranked[1][0]
            if len(ranked) > 2:
                palette.accent = ranked[2][0]
            if len(ranked) > 3:
                palette.additional = [color for color, _ in ranked[3:]]
        except Exception as exc:
            logger.debug("Failed to extract PPTX colors: %s", exc)
        return self._finalize_palette_roles(palette)

    async def _extract_background_color(self, pptx_path: str) -> Optional[str]:
        """从 PPTX 提取背景色（支持 solid fill、gradient 首色、theme color）"""
        try:
            from pptx import Presentation
            prs = Presentation(str(pptx_path))

            for slide in prs.slides:
                bg = slide.background
                fill = bg.fill

                # 尝试多种 fill 类型
                try:
                    fill_type = fill.type
                except Exception:
                    fill_type = None

                if fill_type is not None:
                    try:
                        if hasattr(fill, 'fore_color') and fill.fore_color:
                            rgb = fill.fore_color.rgb
                            if rgb:
                                return f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}"
                    except Exception:
                        pass

                # 尝试从 XML 直接提取背景色
                try:
                    from lxml import etree
                    bg_xml = bg._element
                    # 查找 solidFill / gradFill 中的 srgbClr
                    ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
                    for srgb in bg_xml.iter('{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr'):
                        val = srgb.get('val')
                        if val and len(val) == 6:
                            return f"#{val.lower()}"
                except Exception:
                    pass

            return None
        except Exception as e:
            logger.debug(f"Failed to extract background color: {e}")
            return None

    # ── HTML fallback 方法（当 PPTX 不可用时）──

    def _extract_fonts_from_html(self, slides_html: list[str]) -> Typography:
        typography = Typography()

        # 按 HTML 标签语义分组提取 font-family
        title_tag_re = re.compile(
            r'<(h[1-3])\b[^>]*style\s*=\s*["\'][^"\']*font-family:\s*["\']?([^"\';]+)',
            re.IGNORECASE,
        )
        body_tag_re = re.compile(
            r'<(p|li|span|div)\b[^>]*style\s*=\s*["\'][^"\']*font-family:\s*["\']?([^"\';]+)',
            re.IGNORECASE,
        )
        title_size_re = re.compile(
            r'<h[1-2]\b[^>]*style\s*=\s*["\'][^"\']*font-size:\s*(\d+)',
            re.IGNORECASE,
        )
        body_size_re = re.compile(
            r'<p\b[^>]*style\s*=\s*["\'][^"\']*font-size:\s*(\d+)',
            re.IGNORECASE,
        )

        title_fonts: list[str] = []
        body_fonts: list[str] = []
        title_sizes: list[int] = []
        body_sizes: list[int] = []

        for html in slides_html:
            for m in title_tag_re.finditer(html):
                title_fonts.append(m.group(2).strip().strip("'\""))
            for m in body_tag_re.finditer(html):
                body_fonts.append(m.group(2).strip().strip("'\""))
            title_sizes.extend(int(s) for s in title_size_re.findall(html))
            body_sizes.extend(int(s) for s in body_size_re.findall(html))

        # 如果语义标签没有匹配到字体，fallback 到全局频率
        if not title_fonts and not body_fonts:
            font_pattern = r'font-family:\s*["\']?([^"\';]+)'
            all_fonts = []
            for html in slides_html:
                matches = re.findall(font_pattern, html)
                all_fonts.extend([f.strip().strip("'\"") for f in matches])
            if all_fonts:
                common = Counter(all_fonts).most_common(2)
                if common:
                    typography.title_font = common[0][0]
                    typography.body_font = common[1][0] if len(common) > 1 else common[0][0]
        else:
            if title_fonts:
                typography.title_font = Counter(title_fonts).most_common(1)[0][0]
            if body_fonts:
                typography.body_font = Counter(body_fonts).most_common(1)[0][0]

        if title_sizes:
            title_sizes.sort()
            typography.title_size = title_sizes[len(title_sizes) // 2]
        if body_sizes:
            body_sizes.sort()
            typography.body_size = body_sizes[len(body_sizes) // 2]

        # fallback: 全局统计
        if not title_sizes and not body_sizes:
            size_pattern = r'font-size:\s*(\d+)'
            all_sizes = []
            for html in slides_html:
                all_sizes.extend(int(s) for s in re.findall(size_pattern, html))
            if all_sizes:
                avg_size = sum(all_sizes) / len(all_sizes)
                if avg_size > 24:
                    typography.title_size = int(avg_size)
                else:
                    typography.body_size = int(avg_size)

        return typography

    def _extract_spacing_from_html(self, slides_html: list[str]) -> dict:
        spacing = {}
        for html in slides_html:
            margin_matches = re.findall(r'margin:\s*(\d+)\s*(px|pt)?', html)
            padding_matches = re.findall(r'padding:\s*(\d+)\s*(px|pt)?', html)
            if margin_matches:
                val, unit = margin_matches[0]
                spacing['margin'] = f"{val}{unit or 'pt'}"
            if padding_matches:
                val, unit = padding_matches[0]
                spacing['padding'] = f"{val}{unit or 'pt'}"
            if spacing:
                break
        return spacing



# ═══════════════════════════════════════════════
# ContentPatternExtractor
# ═══════════════════════════════════════════════


class ContentPatternExtractor:
    """从 HTML 和布局中提取内容模式（叙述风格、信息密度、列表风格等）

    LLM 语义分析是核心路径，必须成功才能填充 narrative_style/typical_sections/bullet_style。
    """

    def __init__(self, llm: Optional[Any] = None):
        self.llm = llm

    async def extract(
        self,
        slides_html: list[str],
        slides_text: list[list[str]],
        layout_induction: dict,
        pptx_filename: str = "",
    ) -> ContentPatterns:
        patterns = ContentPatterns()

        # 规则提取：信息密度（基于文本量和元素数）
        info_density = self._estimate_info_density(slides_html, layout_induction)
        if info_density:
            patterns.info_density = info_density

        # LLM 语义分析：narrative_style, typical_sections, bullet_style
        if self.llm:
            semantic_patterns = await self._semantic_analysis(
                slides_html, slides_text, pptx_filename
            )
            if semantic_patterns:
                if semantic_patterns.narrative_style:
                    patterns.narrative_style = semantic_patterns.narrative_style
                if semantic_patterns.info_density:
                    patterns.info_density = semantic_patterns.info_density
                if semantic_patterns.bullet_style:
                    patterns.bullet_style = semantic_patterns.bullet_style
                if semantic_patterns.typical_sections:
                    patterns.typical_sections = semantic_patterns.typical_sections
                if semantic_patterns.max_bullets_per_slide != 5:
                    patterns.max_bullets_per_slide = semantic_patterns.max_bullets_per_slide
        else:
            logger.warning("[TEMPLATE] ContentPatternExtractor: LLM 不可用，narrative_style/typical_sections/bullet_style 将为空")

        return patterns

    def _estimate_info_density(self, slides_html: list[str], layout_induction: dict) -> str:
        """规则估算信息密度"""
        if not slides_html:
            return "medium"
        # 平均每页文本量
        total_text_len = sum(len(re.sub(r'<[^>]+>', '', html)) for html in slides_html)
        avg_text = total_text_len / max(len(slides_html), 1)
        if avg_text > 500:
            return "high"
        elif avg_text < 100:
            return "low"
        return "medium"

    async def _semantic_analysis(
        self,
        slides_html: list[str],
        slides_text: list[list[str]],
        pptx_filename: str,
    ) -> Optional[ContentPatterns]:
        """使用 LLM 进行语义分析。"""
        if not self.llm:
            return None

        llm_callable = await self._get_llm_callable()
        if not llm_callable:
            logger.warning("[TEMPLATE] ContentPatternExtractor: LLM 不可调用，跳过语义分析")
            return None

        # 构建摘要输入（避免 token 过多）
        text_summary = ""
        for i, texts in enumerate(slides_text[:8]):
            page_text = " ".join(texts) if isinstance(texts, list) else str(texts)
            text_summary += f"Slide {i+1}: {page_text[:300]}\n"

        prompt = f"""分析以下 PPT 模板的内容模式。

文件名：{pptx_filename}
幻灯片数量：{len(slides_html)}

各页内容摘要：
{text_summary}

请严格以 JSON 格式输出分析结果（不要添加任何其他文字）：
{{
  "narrative_style": "叙述风格，必须是以下之一：academic/business/technical/educational/creative",
  "info_density": "信息密度，必须是以下之一：high/medium/low",
  "bullet_style": "列表风格，必须是以下之一：numbered/bulleted/icon/mixed/none",
  "typical_sections": ["典型章节名称列表，如：opening, table_of_contents, content, ending"],
  "max_bullets_per_slide": 5
}}
"""
        last_error = None
        for attempt in range(3):
            try:
                response = await llm_callable(prompt)
                if not response:
                    raise ValueError("LLM returned empty response")

                # 解析响应为 dict
                import json as _json
                if isinstance(response, dict):
                    # return_json=True 已经返回了 dict
                    data = response
                elif isinstance(response, str):
                    text = response.strip()
                    # 移除 <think>...</think> 块
                    text = re.sub(r'<think>.*?</think>', '', text, flags=re.DOTALL).strip()
                    # 提取 JSON
                    if "```json" in text:
                        start = text.find("```json") + 7
                        end = text.rfind("```")
                        text = text[start:end].strip()
                    elif "```" in text:
                        start = text.find("```") + 3
                        end = text.rfind("```")
                        text = text[start:end].strip()
                    data = _json.loads(text)
                else:
                    raise ValueError(f"Unexpected response type: {type(response)}")

                result = ContentPatterns()
                result.narrative_style = data.get("narrative_style", "")
                result.info_density = data.get("info_density", "medium")
                result.bullet_style = data.get("bullet_style", "")
                result.typical_sections = data.get("typical_sections", [])
                result.max_bullets_per_slide = int(data.get("max_bullets_per_slide", 5))

                logger.info(
                    f"[TEMPLATE] ContentPatterns LLM 分析成功: "
                    f"narrative={result.narrative_style}, density={result.info_density}, "
                    f"bullet={result.bullet_style}, sections={result.typical_sections}"
                )
                return result

            except Exception as e:
                last_error = e
                logger.warning(f"[TEMPLATE] ContentPatterns LLM 分析 attempt {attempt+1}/3 失败: {e}")

        logger.error(f"[TEMPLATE] ContentPatterns LLM 分析全部失败: {last_error}")
        return None

    async def _get_llm_callable(self):
        """获取可调用的 LLM 函数

        返回的 callable 可能返回 dict（return_json=True）或 str（raw .run()）。
        _semantic_analysis 已经能处理两种类型。
        """
        if self.llm is None:
            return None

        if callable(self.llm):
            logger.info("[TEMPLATE] ContentPatternExtractor: 使用 callable LLM")
            async def _call(prompt):
                try:
                    return await self.llm(prompt, return_json=True)
                except TypeError:
                    return await self.llm(prompt)
            return _call

        if hasattr(self.llm, 'run'):
            logger.info(f"[TEMPLATE] ContentPatternExtractor: 使用 memslides LLM.run() 直接调用")
            async def _call(prompt):
                messages = [{"role": "user", "content": prompt}]
                resp = await self.llm.run(messages)
                if resp and resp.choices:
                    return resp.choices[0].message.content
                return None
            return _call

        logger.warning(f"[TEMPLATE] ContentPatternExtractor: 无法识别的 LLM 类型: {type(self.llm)}")
        return None
